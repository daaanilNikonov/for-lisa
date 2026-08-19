#!/usr/bin/env python3
"""Build «100 к 1» quiz for Клуб продавцов — отработка возражений (кейсы)."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Презентация ГК Форус темный шаблон 16х9 (1).pptx"
OUT_DIR = ROOT / "квиз клуб продавцов"
OUT = OUT_DIR / "Квиз_Клуб_продавцов_Отработка_возражений.pptx"
OUT_COPY = ROOT / "presentation" / "quiz" / "Квиз_Клуб_продавцов_Отработка_возражений.pptx"
TIMER_DIR = ROOT / "presentation" / "quiz"

# Brand
BLUE = RGBColor(0x26, 0xA6, 0xE0)
CARD = RGBColor(0x3F, 0x3F, 0x3F)
CARD_DARK = RGBColor(0x2A, 0x2A, 0x2A)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xBF, 0xBF, 0xBF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GOLD = RGBColor(0xF0, 0xB4, 0x2E)
RED = RGBColor(0xC0, 0x39, 0x2B)

FONT = "Verdana"
POINTS = [10, 30, 50, 70, 100, 150]
L_TITLE = 0
L_CONTENT = 3
L_BG = 6

# Таймеры из файла «Игра клуб продавцов.docx»:
# 10/30/50 → 1–1,5 мин; 70/100 → 2 мин; 150 → 3 мин
def think_seconds(points: int) -> int:
    if points <= 50:
        return 90
    if points <= 100:
        return 120
    return 180


def timer_label(points: int) -> str:
    sec = think_seconds(points)
    if sec % 60 == 0:
        return f"{sec // 60} мин"
    return f"{sec // 60}:{sec % 60:02d}"


# ─── 5 тем × 6 кейсов (из файла, переформулированы в полноценные кейсы) ───
# Однотипные кейсы из файла заменены (отмечены в source).

TOPICS: list[dict] = [
    {
        "name": "Цена и бюджет",
        "short": "Цена",
        "questions": [
            {
                "points": 10,
                "case": (
                    "Сеть из трёх розничных магазинов ищет решение для учёта остатков "
                    "и продаж. Собственник сам на встрече, боль понятная: расхождения "
                    "по складу и ручные своды в Excel. Вы предложили 1С:Розница / УТ "
                    "с небольшим проектом внедрения.\n\n"
                    "Возражение: «Дорого, мы думали до 50 тысяч»."
                ),
                "task": "Отработайте возражение и приведите к следующему шагу покупки.",
                "key": (
                    "Уточнить, из чего сложилась «цифра 50 тыс.»; разложить стоимость "
                    "(лицензии / работы / сопровождение); показать цену ошибки учёта "
                    "на 3 магазинах; предложить MVP на одном магазине или этапность."
                ),
            },
            {
                "points": 30,
                "case": (
                    "Производственная компания просит автоматизировать склад и цех: "
                    "нужны остатки, серии, передача в производство. Есть согласованный "
                    "интерес у начальника производства.\n\n"
                    "Возражение: «В этом году бюджет только на железо, софт — потом»."
                ),
                "task": "Отработайте возражение и приведите к покупке / фиксации сделки.",
                "key": (
                    "Связать железо и софт в один эффект (железо без учёта = полка); "
                    "предложить старт с минимального контура; варианты финансирования "
                    "(этапы, рассрочка, перенос части в OPEX); зафиксировать слот команды "
                    "и дорожную карту «железо + 1С»."
                ),
            },
            {
                "points": 50,
                "case": (
                    "Торговая компания устала от Excel и хочет перейти на 1С:УТ: "
                    "продажи, склад, взаиморасчёты. ЛПР — коммерческий директор, "
                    "сравнил вас с другим интегратором.\n\n"
                    "Возражение: «У конкурента нашли предложение на 40% дешевле»."
                ),
                "task": "Отработайте сравнение и доведите до выбора вашего предложения.",
                "key": (
                    "Сверить состав работ «яблоко к яблоку» (объём, роли, обучение, "
                    "гарантия, SLA); показать риски дешёвого внедрения; предложить "
                    "ценность/этапы, а не войну скидок; дать сравнительную таблицу."
                ),
            },
            {
                "points": 70,
                "case": (
                    # Замена однотипного «скиньте 30% на тендере»
                    "Средний оптовик готов внедрять 1С:УТ, сумма и сроки согласованы. "
                    "На финале финдир меняет модель оплаты.\n\n"
                    "Возражение: «Платим только после запуска в прод, без аванса "
                    "и без оплаты лицензий до go-live — иначе не подписываем»."
                ),
                "task": "Отработайте условие оплаты и приведите к подписанию договора.",
                "key": (
                    "Объяснить, почему лицензии и команда требуют аванса/этапов; "
                    "предложить поэтапную оплату по вехам (обследование / НФТ / запуск); "
                    "ограниченный пилот с понятной предоплатой; не брать 100% постфактум."
                ),
            },
            {
                "points": 100,
                "case": (
                    "Крупный клиент хочет 1С и сопровождение на 3 года: развитие, "
                    "обновления, линия поддержки. Юристы уже в сделке.\n\n"
                    "Возражение: «Оплатим только если гарантируете окупаемость; "
                    "иначе — штраф в договоре»."
                ),
                "task": "Отработайте требование гарантии окупаемости и закройте сделку.",
                "key": (
                    "Отделить ответственность за внедрение от бизнес-ROI клиента; "
                    "предложить KPI проекта (сроки, качество, SLA), а не «штраф за прибыль»; "
                    "совместный расчёт эффекта; пилот с измеримыми метриками."
                ),
            },
            {
                "points": 150,
                "case": (
                    "Финдиректор холдинга готов автоматизировать закрытие месяца "
                    "на базе 1С. Без его подписи сделка не двигается.\n\n"
                    "Возражение: «Покажите ROI за 6 месяцев цифрами, иначе не подпишу»."
                ),
                "task": "Соберите аргументацию ROI и приведите к подписанию.",
                "key": (
                    "Запросить исходные цифры (FTE на закрытие, ошибки, пени, дни soft close); "
                    "построить модель «как есть / как будет»; консервативный сценарий "
                    "на 6–12 мес.; предложить этап с быстрым эффектом (регламент + отчёты)."
                ),
            },
        ],
    },
    {
        "name": "Не сейчас / подумаем",
        "short": "Не сейчас",
        "questions": [
            {
                "points": 10,
                "case": (
                    "ООО хочет электронный документооборот, связанный с 1С: меньше "
                    "бумаги и ручной разноски. Встреча прошла тепло.\n\n"
                    "Возражение: «Напишите КП на почту, мы подумаем»."
                ),
                "task": "Не потеряйте сделку на «КП на почту» — доведите до конкретного шага.",
                "key": (
                    "Уточнить, что именно сравнят и с кем; предложить короткий созвон "
                    "разбора КП; зафиксировать дату и критерий решения; добавить в КП "
                    "следующий шаг (демо/аудит), а не «просто цены»."
                ),
            },
            {
                "points": 30,
                "case": (
                    "Оптовая компания хочет мобильных торговых агентов (заказы с планшета "
                    "в 1С). Сезон высокий, отдел продаж перегружен.\n\n"
                    "Возражение: «Сейчас сезон, вернёмся осенью»."
                ),
                "task": "Свяжите сезон с ценностью и удержите движение к покупке.",
                "key": (
                    "Показать стоимость ручных заказов именно в сезон; предложить "
                    "лёгкий старт (пилотные агенты); забронировать слот внедрения "
                    "на спад с предоплатой/договором сейчас."
                ),
            },
            {
                "points": 50,
                "case": (
                    "После демо 1С:ЗУП HR и главбух «за». Вы ждёте решения.\n\n"
                    "Возражение: «Надо посоветоваться с командой» — и тишина 10 дней."
                ),
                "task": "Верните контакт и доведите до следующего шага покупки.",
                "key": (
                    "Мягкий follow-up с ценностью (не «ну как там?»); выявить, кто "
                    "ещё в круге решения; предложить мини-встречу с возражениями команды; "
                    "дать дедлайн по слоту/акции без давления."
                ),
            },
            {
                "points": 70,
                "case": (
                    "Ритейл хочет онлайн-кассы + связку с 1С. Тема горячая из-за "
                    "требований учёта.\n\n"
                    "Возражение: «Сначала переживём проверку, потом софт»."
                ),
                "task": "Отработайте отсрочку «после проверки» и приведите к старту.",
                "key": (
                    "Показать, что хаос учёта увеличивает риск на проверке; предложить "
                    "минимальный контур «быстрых побед» до/параллельно проверке; "
                    "развести срочное (кассы/чеки) и развитие."
                ),
            },
            {
                "points": 100,
                "case": (
                    # Замена однотипного «заморозьте до конца года»
                    "После согласования цены клиент почти готов к 1С, но меняет процесс.\n\n"
                    "Возражение: «Сначала сделайте бесплатный аудит всех процессов "
                    "на 2 недели — и только по результатам решим, покупать или нет»."
                ),
                "task": "Отработайте требование бесплатного аудита и приведите к оплачиваемому шагу.",
                "key": (
                    "Сузить аудит до 1–2 процессов / 2–3 дней; сделать его платным "
                    "с зачётом в проект; дать осязаемый артефакт (карта болей + оценка); "
                    "не отдавать 2 недели экспертизы «в никуда»."
                ),
            },
            {
                "points": 150,
                "case": (
                    "Стратегический клиент на финале выбора вас как интегратора 1С. "
                    "Устно «всё решили».\n\n"
                    "Возражение: «Старт через 9 месяцев, без предоплаты и без брони "
                    "команды — просто держите нас в уме»."
                ),
                "task": "Зафиксируйте обязательства сторон или честно квалифицируйте отказ.",
                "key": (
                    "Объяснить стоимость «висящего» слота; предложить договор с отложенным "
                    "стартом и небольшой фиксацией; календарный hold; иначе — вернуть "
                    "в pipeline без иллюзии сделки."
                ),
            },
        ],
    },
    {
        "name": "Уже есть / конкурент",
        "short": "Конкуренты",
        "questions": [
            {
                "points": 10,
                "case": (
                    "Складской учёт «на коленке» в Excel тормозит отгрузки: "
                    "ошибки, долгие инвентаризации. Кладовщики против перемен.\n\n"
                    "Возражение: «Excel нас устраивает, привыкли»."
                ),
                "task": "Вскройте боль без атаки на Excel и приведите к демо/пилоту.",
                "key": (
                    "Вопросы про стоимость ошибок и время на сверки; кейс «одного дня "
                    "отгрузки»; предложить пилот на одном складе, не «ломая всё сразу»."
                ),
            },
            {
                "points": 30,
                "case": (
                    "Компания с филиалами хочет единую картину по остаткам и продажам. "
                    "Сейчас у каждого своя локальная логика.\n\n"
                    "Возражение: «У нас уже стоит 1С:Бухгалтерия — этого хватит»."
                ),
                "task": "Покажите границу БП и необходимость контура УТ/ERP — до продажи.",
                "key": (
                    "Отделить регламентированный учёт от операционного; показать, "
                    "чего БП не закрывает (склады, заказы, филиалы); предложить "
                    "надстройку/обмен, а не «выбросить БП»."
                ),
            },
            {
                "points": 50,
                "case": (
                    "Компания хочет CRM + продажи, связку с 1С. Вчера был другой вендор.\n\n"
                    "Возражение: «Нам нарисовали внедрение Битрикс за 2 недели»."
                ),
                "task": "Отработайте «быстрый конкурент» и верните к осознанному выбору.",
                "key": (
                    "Спросить про интеграции с 1С, роли, историю, телефонию, права; "
                    "показать типовые риски «за 2 недели»; предложить сравнительный "
                    "сценарий под их процесс, не антирекламу."
                ),
            },
            {
                "points": 70,
                "case": (
                    "Клиент готов купить лицензии 1С. Деньги на ПО есть, на проект — спорно.\n\n"
                    "Возражение: «Только лицензии, без вашего проекта и обучения»."
                ),
                "task": "Продайте нужный минимум работ или безопасно отпустите в «только лицензии».",
                "key": (
                    "Предупредить о риске «лицензии в стол»; предложить пакет запуска "
                    "(настройка + обучение 2 дней); письменно зафиксировать границы "
                    "ответственности при продаже only-license."
                ),
            },
            {
                "points": 100,
                "case": (
                    "Идёт замена текущего интегратора 1С: качество поддержки упало. "
                    "Формально вас зовут в конкурс.\n\n"
                    "Возражение: «Текущий подрядчик — родственник учредителя; "
                    "формально выбираем вас, фактически — его»."
                ),
                "task": "Определите реальность сделки и найдите путь к покупке или выходу.",
                "key": (
                    "Квалифицировать: есть ли спонсор сильнее родственника; предложить "
                    "пилот/SLA с измеримым эффектом; не демпинговать «для галочки»; "
                    "при политике — вежливый выход с дверью на будущее."
                ),
            },
            {
                "points": 150,
                "case": (
                    "Производство хочет MES-контур на базе 1С (цех, сменно-суточные задания). "
                    "Технически тема живая.\n\n"
                    "Возражение: «Два года назад 1С внедряли — провал, больше не хотим»."
                ),
                "task": "Отработайте травму прошлого внедрения и откройте путь к новой сделке.",
                "key": (
                    "Разобрать, что именно провалилось (цели, подрядчик, данные, люди); "
                    "предложить диагностику «почему снова будет иначе»; маленький пилот "
                    "на одном участке с жёсткими критериями успеха."
                ),
            },
        ],
    },
    {
        "name": "Риски и доверие",
        "short": "Риски",
        "questions": [
            {
                "points": 10,
                "case": (
                    "Клиент рассматривает переход в облачную 1С (1cfresh / частное облако). "
                    "IT насторожен.\n\n"
                    "Возражение: «А вдруг всё упадёт в пятницу вечером?»"
                ),
                "task": "Снимите тревогу и доведите до пилота/договора.",
                "key": (
                    "SLA, резервное копирование, регламент инцидентов; кейсы доступности; "
                    "план отката; предложить пилот на некритичной базе."
                ),
            },
            {
                "points": 30,
                "case": (
                    "Компания хочет КЭДО для кадров (заявления, приказы, подпись), "
                    "интеграция с 1С:ЗУП.\n\n"
                    "Возражение: «Данные сотрудников в облако не отдадим»."
                ),
                "task": "Отработайте ИБ-возражение и предложите приемлемую схему покупки.",
                "key": (
                    "Варианты on-prem / контур заказчика; правовые основания обработки; "
                    "разграничение доступа; пилот на одном юрлице; подключить ИБ рано."
                ),
            },
            {
                "points": 50,
                "case": (
                    "Банк/финконтур + 1С: нужны доработки под внутренние политики. "
                    "КП уже отправлено.\n\n"
                    "Возражение: «ИБ не пропустит без доработок, которых у вас в КП нет»."
                ),
                "task": "Верните сделку в конструктив: что добавить и как продать доработки.",
                "key": (
                    "Запросить чек-лист ИБ; оценить gap-анализ отдельной строкой; "
                    "не обещать «уже всё есть»; этап «требования ИБ» до финальной сметы."
                ),
            },
            {
                "points": 70,
                "case": (
                    "Федеральная сеть хочет централизованную 1С. Сравнивают подрядчиков.\n\n"
                    "Возражение: «Вы работаете только на Иркутск, а Первый Бит — на всю "
                    "Россию: вы не потянете поддержку 24/7»."
                ),
                "task": "Отработайте статус «локальный» и доведите до пилота/контракта.",
                "key": (
                    "Факты: команда, линия, регламент, партнёрская сеть, кейсы удалёнки; "
                    "предложить SLA и эскалации; пилот на регионе; не оправдываться — "
                    "продавать управляемость."
                ),
            },
            {
                "points": 100,
                "case": (
                    "Клиент согласен смотреть 1С, но ставит условие входа.\n\n"
                    "Возражение: «Пилот бесплатно на 3 месяца на всех базах — иначе "
                    "не смотрим»."
                ),
                "task": "Сузьте пилот так, чтобы он продавал, а не разорял — и закройте шаг.",
                "key": (
                    "1 юрлицо / 1 процесс / 2–4 недели; критерии успеха заранее; "
                    "стоимость пилота с зачётом в проект; отказ от «всех баз бесплатно»."
                ),
            },
            {
                "points": 150,
                "case": (
                    "Окологосударственный заказчик хочет 1С через закупку. Юристы жёсткие.\n\n"
                    "Возражение: «В договоре — полная материальная ответственность "
                    "за любой простой и утечку; иначе закупки нет»."
                ),
                "task": "Найдите договорную конструкцию, при которой сделка возможна.",
                "key": (
                    "Ограничение ответственности / страхование / исключения force majeure; "
                    "зона ответственности заказчика (инфра, доступы); эскалация к своему "
                    "юристу; не подписывать безлимит — предложить управляемый риск."
                ),
            },
        ],
    },
    {
        "name": "ЛПР и политика",
        "short": "ЛПР",
        "questions": [
            {
                "points": 10,
                "case": (
                    "IT просит нормальный контур обновлений и поддержки 1С: сейчас "
                    "всё держится на одном внутреннем специалисте.\n\n"
                    "Возражение: «Надо согласовать с закупками, бухгалтерией и охраной "
                    "труда — это на месяцы»."
                ),
                "task": "Превратите «согласования на месяцы» в управляемый план продажи.",
                "key": (
                    "Карта стейкхолдеров и их интересов; короткий one-pager под каждую роль; "
                    "предложить kick-off на 30 минут со всеми; параллелить, не ждать цепочку."
                ),
            },
            {
                "points": 30,
                "case": (
                    "Главбух хочет ускорить отчётность в 1С и явно ваш сторонник.\n\n"
                    "Возражение: «Я не ЛПР, решает директор — и он недоступен»."
                ),
                "task": "Используйте чемпиона и доберите доступ к ЛПР / решение.",
                "key": (
                    "Попросить интро/слот у директора; дать главбуху «скрипт ценности» "
                    "на 5 минут; письмо от вас + от неё; альтернатива — короткий созвон "
                    "в окне директора."
                ),
            },
            {
                "points": 50,
                "case": (
                    "Руководитель продаж хочет воронку и дисциплину сделок в 1С/CRM.\n\n"
                    "Возражение: «IT против любого нового вендора»."
                ),
                "task": "Отработайте блок IT и приведите к совместному решению о покупке.",
                "key": (
                    "Выяснить реальные страхи IT (интеграции, безопасность, нагрузка); "
                    "пригласить IT на демо с их чек-листом; продавать архитектуру и "
                    "сопровождение, не только «хотелки продаж»."
                ),
            },
            {
                "points": 70,
                "case": (
                    "Собственник хочет прозрачную финмодель и управленку на 1С.\n\n"
                    "Возражение: «Финдир и IT тянут в разные решения и блокируют друг друга»."
                ),
                "task": "Сведите стороны к одному решению и движению к договору.",
                "key": (
                    "Фасилитирующая встреча с повесткой критериев выбора; матрица "
                    "требований фин/IT; роль собственника как арбитра; PoC под общие KPI."
                ),
            },
            {
                "points": 100,
                "case": (
                    "По сути вы уже выиграли: вас хвалят, объём понятен. Но всё уходит "
                    "в формальный тендер.\n\n"
                    "Возражение: «По регламенту победит тот, у кого ниже цена в конверте»."
                ),
                "task": "Усильте неценовые критерии и сохраните шанс на победу/сделку.",
                "key": (
                    "Помочь заказчику заложить качество, SLA, опыт, сроки в критерии; "
                    "разбить лоты; ТЗ, где демпинг опасен; параллельно — ценность для ЛПР."
                ),
            },
            {
                "points": 150,
                "case": (
                    "Договор почти подписан, предоплата согласована.\n\n"
                    "Возражение: «Наш внутренний спонсор уволился вчера; новый директор "
                    "всё ставит на паузу»."
                ),
                "task": "За 48 часов спасите сделку или пересоберите цикл продажи.",
                "key": (
                    "Срочно найти нового спонсора/влиятельных союзников; brief для нового "
                    "директора (проблема→эффект→статус); предложить паузу с рамками, "
                    "не исчезновение; мини-демо ценности под его KPI."
                ),
            },
        ],
    },
]


def emu(inches: float) -> int:
    return int(Inches(inches))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_run(run, text, size_pt, bold=False, color=WHITE, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color)
    return box


def add_card(slide, left, top, width, height, fill=CARD, corner=0.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def clear_body_placeholders(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (13, 14, 1, 2):
            if ph.has_text_frame:
                ph.text_frame.clear()


def fill_title(slide, text, size=26):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), text, size, True, WHITE)
            return ph
    return add_textbox(slide, emu(0.97), emu(0.48), emu(11.4), emu(1.0), text, size, True, WHITE)


def fill_shape_text(shape, text, size_pt=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size_pt, bold, color)
    return shape


def link_to_slide(shape, target_slide):
    shape.click_action.target_slide = target_slide


def set_run_scheme_color(run, scheme: str = "hlink"):
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:solidFill") or child.tag.endswith("}solidFill"):
            r_pr.remove(child)
    solid = etree.SubElement(r_pr, qn("a:solidFill"))
    scheme_el = etree.SubElement(solid, qn("a:schemeClr"))
    scheme_el.set("val", scheme)


def fill_shape_text_hyperlink(shape, text, size_pt=18, bold=True):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = FONT
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", FONT)
    set_run_scheme_color(run, "hlink")
    set_run_scheme_color(run, "hlink")
    return shape


def patch_theme_followed_hyperlink_red(prs: Presentation) -> None:
    for part in prs.part.package.iter_parts():
        name = str(getattr(part, "partname", ""))
        if "theme" not in name or not name.endswith(".xml"):
            continue
        try:
            root = etree.fromstring(part.blob)
        except Exception:
            continue
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        changed = False
        hlink = root.find(".//a:clrScheme/a:hlink", ns)
        fol = root.find(".//a:clrScheme/a:folHlink", ns)
        if hlink is not None:
            for child in list(hlink):
                hlink.remove(child)
            srgb = etree.SubElement(hlink, qn("a:srgbClr"))
            srgb.set("val", "FFFFFF")
            changed = True
        if fol is not None:
            for child in list(fol):
                fol.remove(child)
            srgb = etree.SubElement(fol, qn("a:srgbClr"))
            srgb.set("val", "C0392B")
            changed = True
        if changed:
            part._blob = etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            )


def bind_text_run_hyperlink(shape) -> None:
    sp = shape._element
    hlink = None
    for el in sp.iter():
        if el.tag == qn("a:hlinkClick"):
            hlink = el
            break
    if hlink is None:
        return
    r_id = hlink.get(qn("r:id"))
    action = hlink.get("action")
    if not r_id:
        return
    for r in sp.iter(qn("a:r")):
        r_pr = r.find(qn("a:rPr"))
        if r_pr is None:
            r_pr = etree.Element(qn("a:rPr"))
            r.insert(0, r_pr)
        for old in list(r_pr):
            if old.tag == qn("a:hlinkClick") or old.tag.endswith("}hlinkClick"):
                r_pr.remove(old)
        for child in list(r_pr):
            if child.tag == qn("a:solidFill") or child.tag.endswith("}solidFill"):
                r_pr.remove(child)
        solid = etree.SubElement(r_pr, qn("a:solidFill"))
        scheme_el = etree.SubElement(solid, qn("a:schemeClr"))
        scheme_el.set("val", "hlink")
        r_pr.set("u", "none")
        hl = etree.SubElement(r_pr, qn("a:hlinkClick"))
        hl.set(qn("r:id"), r_id)
        if action:
            hl.set("action", action)


def fill_shape_key(shape, key_text: str):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.TOP)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    set_run(p0.add_run(), "Ключ для ведущего", 11, True, BLUE)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(4)
    p1.alignment = PP_ALIGN.LEFT
    set_run(p1.add_run(), key_text, 11, False, NEAR_BLACK)
    try:
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
    except Exception:
        pass
    return shape


def add_appear_after_ms(slide, shape, delay_ms: int) -> None:
    spid = str(shape.shape_id)
    sld = slide._element
    for old in list(sld.findall(qn("p:timing"))):
        sld.remove(old)
    timing_xml = f"""
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="{delay_ms}"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:par>
                                  <p:cTn id="5" presetID="1" presetClass="entr" presetSubtype="0"
                                         fill="hold" grpId="0" nodeType="withEffect">
                                    <p:stCondLst>
                                      <p:cond delay="0"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                      <p:set>
                                        <p:cBhvr>
                                          <p:cTn id="6" dur="1" fill="hold">
                                            <p:stCondLst>
                                              <p:cond delay="0"/>
                                            </p:stCondLst>
                                          </p:cTn>
                                          <p:tgtEl>
                                            <p:spTgt spid="{spid}"/>
                                          </p:tgtEl>
                                          <p:attrNameLst>
                                            <p:attrName>style.visibility</p:attrName>
                                          </p:attrNameLst>
                                        </p:cBhvr>
                                        <p:to>
                                          <p:strVal val="visible"/>
                                        </p:to>
                                      </p:set>
                                    </p:childTnLst>
                                  </p:cTn>
                                </p:par>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
      <p:bldLst>
        <p:bldP spid="{spid}" grpId="0" animBg="1"/>
      </p:bldLst>
    </p:timing>
    """
    sld.append(etree.fromstring(timing_xml))


def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "Клуб продавцов", 20, False, SOFT)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(8)
            set_run(p2.add_run(), "Отработка возражений", 30, True, BLUE)
            p3 = tf.add_paragraph()
            p3.space_before = Pt(14)
            set_run(p3.add_run(), "Формат «100 к 1»  ·  командная Своя игра", 15, True, WHITE)
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "5 тем  ·  30 кейсов  ·  10 / 30 / 50 / 70 / 100 / 150\n"
                "Таймер: 1,5 мин (10–50)  ·  2 мин (70–100)  ·  3 мин (150)\n"
                "При неверном ответе вопрос может перейти другой команде",
                13,
                False,
                SOFT,
            )
    return slide


def build_board(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Игровое поле  ·  100 к 1", 24)
    clear_body_placeholders(slide)
    add_textbox(
        slide,
        emu(0.7),
        emu(1.15),
        emu(11.8),
        emu(0.35),
        "Клик по баллам → кейс  ·  открытые ячейки краснеют  ·  таймер по сложности  ·  «К полю» — назад",
        11,
        False,
        SOFT,
    )

    topic_w = emu(2.35)
    cell_w = emu(1.45)
    cell_h = emu(0.78)
    gap_x = emu(0.1)
    gap_y = emu(0.1)
    left0 = emu(0.55)
    top0 = emu(1.6)

    theme_hdr = add_card(slide, left0, top0, topic_w, cell_h, CARD_DARK, 0.08)
    fill_shape_text(theme_hdr, "Тема", 13, True, SOFT)
    for ci, pts in enumerate(POINTS):
        left = left0 + topic_w + gap_x + ci * (cell_w + gap_x)
        hdr = add_card(slide, left, top0, cell_w, cell_h, BLUE, 0.1)
        fill_shape_text(hdr, str(pts), 16, True, WHITE)

    cell_shapes: dict[tuple[int, int], object] = {}
    for ti, topic in enumerate(TOPICS):
        top = top0 + (ti + 1) * (cell_h + gap_y)
        topic_card = add_card(slide, left0, top, topic_w, cell_h, CARD, 0.08)
        fill_shape_text(topic_card, topic["name"], 11, True, WHITE)
        for qi, q in enumerate(topic["questions"]):
            left = left0 + topic_w + gap_x + qi * (cell_w + gap_x)
            fill = GOLD if q["points"] >= 100 else BLUE
            shape = add_card(slide, left, top, cell_w, cell_h, fill, 0.12)
            fill_shape_text_hyperlink(shape, str(q["points"]), 18, True)
            cell_shapes[(ti, qi)] = shape
    return slide, cell_shapes


def build_case_slide(prs, topic, qdata, board_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, topic["name"], 20)
    clear_body_placeholders(slide)

    pts = qdata["points"]
    sec = think_seconds(pts)
    delay_ms = sec * 1000
    gif = TIMER_DIR / f"timer_{sec}s.gif"

    badge = add_card(slide, emu(10.55), emu(0.38), emu(1.55), emu(0.5), GOLD, 0.2)
    fill_shape_text(badge, f"{pts}", 16, True, NEAR_BLACK)

    if gif.exists():
        slide.shapes.add_picture(str(gif), emu(10.4), emu(1.0), width=emu(1.85))
    else:
        tcard = add_card(slide, emu(10.4), emu(1.0), emu(1.85), emu(1.85), CARD, 0.15)
        fill_shape_text(tcard, str(sec), 28, True, GOLD)

    add_textbox(
        slide,
        emu(10.3),
        emu(2.9),
        emu(2.05),
        emu(0.45),
        f"Таймер {timer_label(pts)}",
        11,
        True,
        SOFT,
        PP_ALIGN.CENTER,
    )

    # Case card
    case_card = add_card(slide, emu(0.55), emu(1.2), emu(9.6), emu(3.35), CARD, 0.08)
    tf = case_card.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.TOP)
    try:
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
    except Exception:
        pass
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    set_run(p0.add_run(), "Кейс", 11, True, GOLD)
    # Split case into paragraphs for readability
    for block in qdata["case"].split("\n\n"):
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        size = 13 if len(qdata["case"]) < 420 else 12
        set_run(p.add_run(), block, size, False, WHITE)

    # Task line
    add_textbox(
        slide,
        emu(0.7),
        emu(4.65),
        emu(9.3),
        emu(0.45),
        "Задача: " + qdata["task"],
        12,
        True,
        GOLD,
        PP_ALIGN.LEFT,
    )

    # Key card — appears after timer
    key_card = add_card(slide, emu(0.55), emu(5.15), emu(9.6), emu(1.45), CARD_LIGHT, 0.08)
    fill_shape_key(key_card, qdata["key"])
    add_appear_after_ms(slide, key_card, delay_ms=delay_ms)

    add_textbox(
        slide,
        emu(10.3),
        emu(3.4),
        emu(2.05),
        emu(0.85),
        f"Ключ ведущего\nчерез {timer_label(pts)}",
        10,
        False,
        GOLD,
        PP_ALIGN.CENTER,
    )

    back = add_card(slide, emu(10.4), emu(5.45), emu(1.85), emu(1.05), BLUE, 0.12)
    fill_shape_text(back, "← К полю", 13, True, WHITE)
    link_to_slide(back, board_slide)
    return slide


def sort_topic_questions():
    for topic in TOPICS:
        topic["questions"].sort(key=lambda q: POINTS.index(q["points"]))


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    for sec in (90, 120, 180):
        gif = TIMER_DIR / f"timer_{sec}s.gif"
        if not gif.exists():
            raise SystemExit(f"Timer GIF not found: {gif}")

    sort_topic_questions()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    OUT_COPY.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    delete_all_slides(prs)

    build_title(prs)
    board_slide, cell_shapes = build_board(prs)

    for ti, topic in enumerate(TOPICS):
        for qi, qdata in enumerate(topic["questions"]):
            q_slide = build_case_slide(prs, topic, qdata, board_slide)
            cell = cell_shapes.get((ti, qi))
            if cell is not None:
                link_to_slide(cell, q_slide)
                bind_text_run_hyperlink(cell)
                try:
                    cell.name = f"Cell_{ti}_{qi}"
                except Exception:
                    pass

    patch_theme_followed_hyperlink_red(prs)
    prs.save(str(OUT))
    shutil.copy2(OUT, OUT_COPY)

    n_q = sum(len(t["questions"]) for t in TOPICS)
    print(f"Saved: {OUT}")
    print(f"Copy:  {OUT_COPY}")
    print(f"Slides: {len(prs.slides)} (title + board + {n_q} cases)")
    print("Timers: 90s (10–50), 120s (70–100), 180s (150)")
    print(f"Topics: {', '.join(t['name'] for t in TOPICS)}")


if __name__ == "__main__":
    main()
