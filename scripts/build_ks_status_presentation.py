#!/usr/bin/env python3
"""Жёлто-белая презентация Форус: статус работы с «1С:Кабинет сотрудника»."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "Статус_Кабинет_сотрудника_август_2026.pptx"
BRAND = ROOT / "presentation" / "assets" / "brand"

# Форус · шафран / белый (не тёмно-синий шаблон)
YELLOW = RGBColor(0xFE, 0xCF, 0x68)
YELLOW_DEEP = RGBColor(0xE8, 0xB4, 0x3A)
CREAM = RGBColor(0xFF, 0xF8, 0xE8)
CREAM2 = RGBColor(0xFF, 0xF1, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SLATE = RGBColor(0x39, 0x5F, 0x75)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
LINE = RGBColor(0xF0, 0xE0, 0xB8)
SOFT_BG = RGBColor(0xFF, 0xFC, 0xF6)

FONT = "Verdana"
SW, SH = Inches(13.333), Inches(7.5)


def emu(inches: float) -> int:
    return int(Inches(inches))


def set_run(run, text, size_pt, bold=False, color=INK, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color)
    return box


def add_multitext(slide, left, top, width, height, paragraphs, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of (text, size, bold, color)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    for i, item in enumerate(paragraphs):
        text, size, bold, color = item[:4]
        space_before = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(space_before)
        set_run(p.add_run(), text, size, bold, color)
    return box


def add_rect(slide, left, top, width, height, fill, corner=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if corner is not None:
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    return shape


def add_card(slide, left, top, width, height, fill=CREAM, corner=0.08):
    return add_rect(slide, left, top, width, height, fill, corner)


def add_picture(slide, name, left, top, width=None, height=None):
    path = BRAND / name
    if not path.exists():
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, 0, 0, SW, SH, WHITE)
    return slide


def add_footer(slide, page: int, total: int):
    add_rect(slide, 0, emu(7.28), SW, emu(0.22), YELLOW)
    add_textbox(
        slide, emu(0.5), emu(7.08), emu(8.5), emu(0.22),
        "ГК Форус  ·  группа продуктового запуска  ·  1С:Кабинет сотрудника",
        9, False, MUTED,
    )
    add_textbox(
        slide, emu(11.4), emu(7.08), emu(1.4), emu(0.22),
        f"{page} / {total}",
        9, False, MUTED, PP_ALIGN.RIGHT,
    )


def add_header(slide, kicker: str, title: str):
    add_picture(slide, "logo-forus.png", emu(0.5), emu(0.22), height=emu(0.38))
    add_picture(slide, "wave_tr.png", emu(10.55), emu(-0.05), width=emu(2.9))
    add_textbox(slide, emu(0.5), emu(0.72), emu(12.2), emu(0.28), kicker, 11, True, SLATE)
    add_textbox(slide, emu(0.5), emu(0.98), emu(12.2), emu(0.55), title, 26, True, INK)
    add_rect(slide, emu(0.5), emu(1.55), emu(1.35), emu(0.07), YELLOW)


def numbered_dot(slide, left, top, num, size=0.42):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, emu(size), emu(size))
    circ.fill.solid()
    circ.fill.fore_color.rgb = YELLOW
    circ.line.fill.background()
    add_textbox(
        slide, left, top, emu(size), emu(size),
        str(num), 11, True, INK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
    )
    return circ


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
TOTAL = 12


def slide_title(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, emu(0.22), SH, YELLOW)
    add_picture(s, "logo-forus.png", emu(0.7), emu(0.45), height=emu(0.55))
    add_picture(s, "wave_tr.png", emu(9.8), emu(-0.1), width=emu(3.7))
    add_textbox(s, emu(0.7), emu(1.7), emu(11.5), emu(0.35),
                "Группа продуктового запуска  ·  август 2026", 14, True, SLATE)
    add_textbox(s, emu(0.7), emu(2.15), emu(12.0), emu(1.5),
                "Работа с сервисом\n«1С:Кабинет сотрудника»", 36, True, INK)
    add_rect(s, emu(0.7), emu(3.85), emu(2.0), emu(0.1), YELLOW)
    add_textbox(
        s, emu(0.7), emu(4.2), emu(11.0), emu(1.2),
        "Статус проекта: карта работы, готовые материалы,\n"
        "гипотезы в проверке и фактическая воронка по базам.",
        16, False, MUTED,
    )
    add_card(s, emu(0.7), emu(5.7), emu(3.5), emu(1.05), CREAM, 0.1)
    add_textbox(s, emu(0.9), emu(5.82), emu(3.2), emu(0.35), "Период факта", 11, False, SLATE)
    add_textbox(s, emu(0.9), emu(6.12), emu(3.2), emu(0.45), "23.07 — 24.08.2026", 16, True, INK)
    add_card(s, emu(4.4), emu(5.7), emu(3.7), emu(1.05), CREAM, 0.1)
    add_textbox(s, emu(4.6), emu(5.82), emu(3.4), emu(0.35), "Команда на линии", 11, False, SLATE)
    add_textbox(s, emu(4.6), emu(6.12), emu(3.4), emu(0.45), "Юлиана · Соня · Данил", 16, True, INK)
    add_card(s, emu(8.3), emu(5.7), emu(4.3), emu(1.05), YELLOW, 0.1)
    add_textbox(s, emu(8.5), emu(5.82), emu(4.0), emu(0.35), "Прозвонено за период", 11, False, INK)
    add_textbox(s, emu(8.5), emu(6.12), emu(4.0), emu(0.45), "1 084 звонка", 16, True, INK)
    add_picture(s, "wave_bl.png", emu(-0.15), emu(6.55), width=emu(3.2))
    return s


def slide_agenda(prs):
    s = blank_slide(prs)
    add_header(s, "Содержание", "О чём этот отчёт")
    items = [
        ("01", "Проект", "Зачем запускаем Кабинет сотрудника и как устроена группа"),
        ("02", "Карта проекта", "Четыре контура работы: свои, холод, семинары, связка продуктов"),
        ("03", "Материалы", "Что уже собрано: КП, скрипты, регламент, модель, аналитика"),
        ("04", "Гипотезы", "Четыре проверяемых ставки — формулировка, цель, как смотрим факт"),
        ("05", "Воронка", "Звонки, успешные, сделки, записи и демо по базам"),
        ("06", "Кейсы демо", "Место под результаты проведённых демонстраций"),
    ]
    top = emu(1.85)
    for i, (num, title, body) in enumerate(items):
        y = top + emu(i * 0.8)
        add_card(s, emu(0.5), y, emu(12.3), emu(0.72), CREAM if i % 2 == 0 else SOFT_BG, 0.08)
        add_textbox(s, emu(0.7), y + emu(0.12), emu(0.8), emu(0.5), num, 20, True, YELLOW_DEEP,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(s, emu(1.6), y + emu(0.08), emu(3.2), emu(0.55), title, 16, True, INK,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(s, emu(5.0), y + emu(0.08), emu(7.5), emu(0.55), body, 13, False, MUTED,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_footer(s, 2, TOTAL)
    return s


def slide_project(prs):
    s = blank_slide(prs)
    add_header(s, "1. Описание проекта", "Сервис, который закрывает рутину кадров и бухгалтерии")
    add_card(s, emu(0.5), emu(1.85), emu(7.6), emu(4.85), CREAM, 0.08)
    add_textbox(s, emu(0.75), emu(2.05), emu(7.1), emu(0.4), "Что продаём", 13, True, SLATE)
    add_textbox(
        s, emu(0.75), emu(2.45), emu(7.1), emu(1.55),
        "«1С:Кабинет сотрудника» — облачный КЭДО в экосистеме 1С: "
        "заявления, отпуска, расчётные листки по ст. 136 ТК РФ, подпись, "
        "согласования и мобильный доступ. Клиенту не нужно менять учётную систему.",
        14, False, INK,
    )
    bullets = [
        ("Для клиента", "Меньше бумаги и ручных заявок, прозрачный кадровый контур, юридически значимый обмен."),
        ("Для Форус", "Допродажа в действующую базу ИТС/ЗУП и вход в новые компании через понятный сервис."),
        ("Почему сейчас", "Рынок уже требует КЭДО, но порог входа высокий — мы снимаем его демо и настройкой."),
    ]
    y = emu(4.1)
    for title, body in bullets:
        add_rect(s, emu(0.75), y + emu(0.08), emu(0.14), emu(0.14), YELLOW)
        add_textbox(s, emu(1.05), y, emu(6.8), emu(0.28), title, 13, True, INK)
        add_textbox(s, emu(1.05), y + emu(0.28), emu(6.8), emu(0.5), body, 12, False, MUTED)
        y += emu(0.78)

    add_card(s, emu(8.3), emu(1.85), emu(4.5), emu(4.85), INK, 0.08)
    add_textbox(s, emu(8.55), emu(2.1), emu(4.05), emu(0.4), "Группа запуска", 13, True, YELLOW)
    add_textbox(
        s, emu(8.55), emu(2.55), emu(4.05), emu(1.3),
        "Не «просто отдел продаж», а контур проверки гипотез: "
        "кто берёт сервис, с каким оффером и из какой базы.",
        13, False, WHITE,
    )
    kpis = [
        ("3", "менеджера на линии"),
        ("4", "гипотезы в работе"),
        ("2 300+", "ИНН целевой холодной базы по ОКВЭД"),
        ("1 084", "звонка в факте за месяц"),
    ]
    y = emu(4.0)
    for n, label in kpis:
        add_textbox(s, emu(8.55), y, emu(1.5), emu(0.4), n, 18, True, YELLOW)
        add_textbox(s, emu(10.1), y, emu(2.4), emu(0.45), label, 12, False, RGBColor(0xEE, 0xEE, 0xEE),
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        y += emu(0.55)
    add_footer(s, 3, TOTAL)
    return s


def slide_map(prs):
    s = blank_slide(prs)
    add_header(s, "1. Карта проекта", "Как устроена ежедневная работа группы")
    add_textbox(
        s, emu(0.5), emu(1.72), emu(12.3), emu(0.45),
        "Карта дня менеджеров + четыре контура касания. Воронка общая: "
        "прозвон → успешный контакт → сделка → запись → демо → продажа.",
        13, False, MUTED,
    )
    contours = [
        ("A", "Свои клиенты", "ЗУП, КП, канбан", "Оффер: «это уже входит в тариф ИТС ПРОФ» — демо без барьера цены."),
        ("B", "Холод · ОКВЭД", "вахта, общепит, производство", "Берём только фирмы с подходящим видом деятельности из модели ICP."),
        ("C", "Семинары", "вебинары, заявки с сайта", "Тёплый след: отвечаем на вопросы, закрепляем пользу, ведём в демо."),
        ("D", "Связка продуктов", "ДОКИ Логистика + КЭДО", "Кабинет как продолжение документооборота, а не отдельная «ещё одна подписка»."),
    ]
    for i, (lit, title, bases, body) in enumerate(contours):
        left = emu(0.5) + emu(i * 3.2)
        add_card(s, left, emu(2.3), emu(3.05), emu(3.55), CREAM, 0.08)
        add_rect(s, left, emu(2.3), emu(3.05), emu(0.12), YELLOW)
        add_textbox(s, left + emu(0.18), emu(2.55), emu(2.7), emu(0.35), lit, 18, True, YELLOW_DEEP)
        add_textbox(s, left + emu(0.18), emu(2.95), emu(2.7), emu(0.7), title, 16, True, INK)
        add_textbox(s, left + emu(0.18), emu(3.6), emu(2.7), emu(0.55), bases, 11, True, SLATE)
        add_textbox(s, left + emu(0.18), emu(4.2), emu(2.7), emu(1.4), body, 12, False, MUTED)
    add_card(s, emu(0.5), emu(6.05), emu(12.3), emu(0.85), SOFT_BG, 0.08)
    add_textbox(
        s, emu(0.7), emu(6.18), emu(12.0), emu(0.6),
        "Операционный слой: «Карта дня» (чеклист и доска каждого менеджера), "
        "регламент после записи на демо, свод «КС аналитика» по трём листам.",
        13, False, INK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, 4, TOTAL)
    return s


def slide_materials(prs):
    s = blank_slide(prs)
    add_header(s, "2. Что уже сделано", "Пакет материалов, на котором едет линия")
    cols = [
        (
            "Клиенту в руки",
            [
                "Листовка «Кабинет сотрудника»",
                "КП на внедрение КЭДО (только КС)",
                "КП-пакет 750 кабинетов + акция",
                "Преддемонстрация для МПП",
                "Памятка после оплаты пакета",
                "Листовка акции «Больше, чем кешбэк»",
            ],
        ),
        (
            "Менеджеру на линию",
            [
                "Скрипт КабС / ДОКИ / Смартвей",
                "Регламент после записи на демо",
                "Мастерская по КС (июнь 2026)",
                "Квиз 1С-ЭПД · Доки · Логистика",
                "Карта дня группы запуска",
                "Блок-схема воронки работы с клиентом",
            ],
        ),
        (
            "Управление и экономика",
            [
                "Свод «КС аналитика» (3 менеджера)",
                "Инвестиционная модель и план 3/5/7",
                "Целевая база ОКВЭД: 2 306 ИНН",
                "Экономика перехода на КЭДО",
                "Прайсы и линейка кабинетов",
                "Пилот «благодарность + тест-драйв»",
            ],
        ),
    ]
    for i, (title, items) in enumerate(cols):
        left = emu(0.5) + emu(i * 4.2)
        add_card(s, left, emu(1.85), emu(4.0), emu(5.05), CREAM, 0.08)
        add_rect(s, left, emu(1.85), emu(4.0), emu(0.7), YELLOW)
        add_textbox(s, left + emu(0.2), emu(1.95), emu(3.6), emu(0.5), title, 15, True, INK,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        y = emu(2.75)
        for n, item in enumerate(items, 1):
            numbered_dot(s, left + emu(0.22), y, n, 0.32)
            add_textbox(s, left + emu(0.65), y - emu(0.02), emu(3.15), emu(0.4), item, 12, False, INK,
                        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            y += emu(0.62)
    add_footer(s, 5, TOTAL)
    return s


def slide_hypotheses_overview(prs):
    s = blank_slide(prs)
    add_header(s, "3. Гипотезы", "Четыре ставки, которые сейчас проверяем звонками")
    hyps = [
        ("H1", "Свои · ИТС ПРОФ", "+15%", "к конверсии в запись",
         "Если предлагаем демо формулировкой «у вас это входит в тариф ИТС ПРОФ», запись растёт."),
        ("H2", "Холод · ОКВЭД", "выше", "конверсия в сделку",
         "Если берём только фирмы с целевыми ОКВЭД (торговля, производство, вахта, общепит…), сделки идут чаще."),
        ("H3", "Холод · связка", "до 15%", "конверсия в сделку",
         "Если продаём ДОКИ Логистика + КЭДО Кабинет сотрудника пакетом, вероятность сделки выше."),
        ("H4", "Семинары", "+10%", "к конверсии vs первый обзвон",
         "Если после семинара даём пользу и закрываем вопросы — клиент «залояливается», конверсия выше холодного первого звонка."),
    ]
    for i, (code, seg, metric, metric_l, body) in enumerate(hyps):
        col, row = i % 2, i // 2
        left = emu(0.5) + emu(col * 6.4)
        top = emu(1.85) + emu(row * 2.5)
        add_card(s, left, top, emu(6.2), emu(2.3), CREAM, 0.08)
        add_rect(s, left, top, emu(0.14), emu(2.3), YELLOW)
        add_textbox(s, left + emu(0.4), top + emu(0.18), emu(1.2), emu(0.4), code, 16, True, YELLOW_DEEP)
        add_textbox(s, left + emu(1.6), top + emu(0.2), emu(4.3), emu(0.38), seg, 14, True, SLATE)
        add_textbox(s, left + emu(0.4), top + emu(0.65), emu(2.4), emu(0.5), metric, 26, True, INK)
        add_textbox(s, left + emu(2.8), top + emu(0.75), emu(3.1), emu(0.4), metric_l, 12, False, MUTED,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(s, left + emu(0.4), top + emu(1.25), emu(5.5), emu(0.85), body, 13, False, INK)
    add_footer(s, 6, TOTAL)
    return s


def slide_h_warm(prs):
    s = blank_slide(prs)
    add_header(s, "3. Гипотезы · свои клиенты", "H1 ИТС ПРОФ и H4 база семинаров")
    # H1
    add_card(s, emu(0.5), emu(1.85), emu(12.3), emu(2.45), CREAM, 0.08)
    add_textbox(s, emu(0.75), emu(2.0), emu(1.0), emu(0.35), "H1", 16, True, YELLOW_DEEP)
    add_textbox(s, emu(1.7), emu(2.02), emu(10.8), emu(0.35),
                "База наших клиентов  ·  «входит в тариф ИТС ПРОФ»", 16, True, INK)
    add_textbox(
        s, emu(0.75), emu(2.45), emu(12.0), emu(0.7),
        "Если предлагаем демо Кабинета сотрудника через причину «у вас это уже входит в тариф ИТС ПРОФ», "
        "конверсия в запись на демо растёт на 15% относительно обычного обзвона своих.",
        13, False, INK,
    )
    facts = [
        ("Как проверяем", "ЗУП + КП + канбан — контур своих. Смотрим долю записей от успешных и от сделок."),
        ("Факт периода", "227 прозвонов · 27 в сделку (11,9%) · 10 записей · 8 демо. Запись от сделок: 37%."),
        ("Цель H1", "+15% к конверсии в запись. Сравниваем с первым касанием без оффера ИТС."),
    ]
    x = emu(0.75)
    for t, b in facts:
        add_textbox(s, x, emu(3.2), emu(3.9), emu(0.28), t, 11, True, SLATE)
        add_textbox(s, x, emu(3.48), emu(3.9), emu(0.65), b, 11, False, MUTED)
        x += emu(4.0)

    add_card(s, emu(0.5), emu(4.5), emu(12.3), emu(2.35), SOFT_BG, 0.08)
    add_textbox(s, emu(0.75), emu(4.65), emu(1.0), emu(0.35), "H4", 16, True, YELLOW_DEEP)
    add_textbox(s, emu(1.7), emu(4.67), emu(10.8), emu(0.35),
                "База наших клиентов с семинаров  ·  польза и ответы на вопросы", 16, True, INK)
    add_textbox(
        s, emu(0.75), emu(5.1), emu(12.0), emu(0.65),
        "Если после семинара/вебинара даём пользу и закрываем вопросы, клиент закрепляется за Форус, "
        "а конверсия выше на 10%, чем при первом обзвоне холодной или «неразогретой» своей базы.",
        13, False, INK,
    )
    facts2 = [
        ("Как проверяем", "База вебинаров + заявки с сайта vs первый обзвон прочих контуров."),
        ("Факт периода", "302 прозвона · 17 в сделку (5,6%) · 1 запись · 1 демо. Контакт тёплый, запись пока редкая."),
        ("Цель H4", "Конверсия на 10% выше, чем при первом обзвоне. Нужен хвост после демо, чтобы честно сравнить."),
    ]
    x = emu(0.75)
    for t, b in facts2:
        add_textbox(s, x, emu(5.8), emu(3.9), emu(0.28), t, 11, True, SLATE)
        add_textbox(s, x, emu(6.08), emu(3.9), emu(0.6), b, 11, False, MUTED)
        x += emu(4.0)
    add_footer(s, 7, TOTAL)
    return s


def slide_h_cold(prs):
    s = blank_slide(prs)
    add_header(s, "3. Гипотезы · холодная база", "H2 целевые ОКВЭД и H3 связка ДОКИ + КЭДО")
    add_card(s, emu(0.5), emu(1.85), emu(12.3), emu(2.5), CREAM, 0.08)
    add_textbox(s, emu(0.75), emu(2.0), emu(1.0), emu(0.35), "H2", 16, True, YELLOW_DEEP)
    add_textbox(s, emu(1.7), emu(2.02), emu(10.8), emu(0.35),
                "Холодная база  ·  только компании с подходящими ОКВЭД", 16, True, INK)
    add_textbox(
        s, emu(0.75), emu(2.45), emu(12.0), emu(0.55),
        "Если продаём Кабинет сотрудника фирмам из ICP-среза, конверсия в сделку выше, "
        "чем по «любому холодку». Срез модели: 2 306 ИНН, горячих (≥22 чел.) — 349.",
        13, False, INK,
    )
    add_textbox(
        s, emu(0.75), emu(3.05), emu(12.0), emu(0.45),
        "ОКВЭД: торговля 45–47 · производство 10–33 · стройка 41–43 · транспорт/склад 49, 52 · "
        "HoReCa 55–56 · недвижимость 68 · кадры 78 · охрана/админ. 80–82 · финансы, IT, образование, медицина.",
        12, False, SLATE,
    )
    add_textbox(
        s, emu(0.75), emu(3.55), emu(12.0), emu(0.55),
        "Факт контура (вахта + общепит + холодка общепит + производство + новая КС): "
        "553 прозвона · 25 в сделку (4,5%) · 1 запись · 2 демо. Сделки есть, запись на демо — узкое место.",
        12, False, MUTED,
    )

    add_card(s, emu(0.5), emu(4.55), emu(12.3), emu(2.3), SOFT_BG, 0.08)
    add_textbox(s, emu(0.75), emu(4.7), emu(1.0), emu(0.35), "H3", 16, True, YELLOW_DEEP)
    add_textbox(s, emu(1.7), emu(4.72), emu(10.8), emu(0.35),
                "Холодная база  ·  связка «ДОКИ Логистика + КЭДО Кабинет сотрудника»", 16, True, INK)
    add_textbox(
        s, emu(0.75), emu(5.15), emu(12.0), emu(0.7),
        "Если предлагаем клиенту не один сервис, а связку документооборота логистики и кадрового КЭДО, "
        "вероятность довести до сделки растёт, цель — конверсия в сделку до 15%.",
        13, False, INK,
    )
    add_textbox(
        s, emu(0.75), emu(5.9), emu(12.0), emu(0.7),
        "Как проверяем: оффер связки на холодных с логистическим/производственным профилем "
        "и на переданных контактах. Скрипт КабС+ДОКИ уже есть. В факте периода отдельный тег связки "
        "ещё копится — фиксируем в аналитике как отдельный источник.",
        12, False, MUTED,
    )
    add_footer(s, 8, TOTAL)
    return s


def slide_funnel_totals(prs):
    s = blank_slide(prs)
    add_header(s, "4. Показатели по базам", "Воронка за 23.07–24.08.2026 · три менеджера")
    kpis = [
        ("1 084", "прозвонено", "100%"),
        ("483", "успешных", "44,6%"),
        ("70", "в сделку", "6,5%"),
        ("12", "записей", "1,1%"),
        ("11", "демо провели", "1,0%"),
        ("0", "продаж", "пока 0"),
    ]
    for i, (n, label, pct) in enumerate(kpis):
        left = emu(0.5) + emu(i * 2.12)
        fill = YELLOW if i == 0 else CREAM
        add_card(s, left, emu(1.85), emu(2.02), emu(1.85), fill, 0.1)
        add_textbox(s, left + emu(0.1), emu(2.0), emu(1.82), emu(0.7), n, 26, True, INK,
                    PP_ALIGN.CENTER)
        add_textbox(s, left + emu(0.1), emu(2.7), emu(1.82), emu(0.4), label, 12, True, SLATE,
                    PP_ALIGN.CENTER)
        add_textbox(s, left + emu(0.1), emu(3.1), emu(1.82), emu(0.35), pct, 12, False, MUTED,
                    PP_ALIGN.CENTER)

    add_textbox(s, emu(0.5), emu(3.9), emu(12.3), emu(0.35), "Кто набрал цифру", 14, True, INK)
    managers = [
        ("Юлиана Юнусова", "114", "60", "24", "9", "7"),
        ("Соня Оглоблина", "555", "238", "29", "2", "3"),
        ("Данил Кургузов", "415", "185", "17", "1", "1"),
    ]
    headers = ["Менеджер", "Прозвонено", "Успешных", "В сделку", "Записи", "Демо"]
    # header row
    y = emu(4.35)
    add_rect(s, emu(0.5), y, emu(12.3), emu(0.42), INK)
    widths = [3.3, 1.8, 1.8, 1.8, 1.8, 1.8]
    x = emu(0.5)
    for w, h in zip(widths, headers):
        add_textbox(s, x, y, emu(w), emu(0.42), h, 12, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        x += emu(w)
    for i, row in enumerate(managers):
        y = emu(4.77) + emu(i * 0.55)
        bg = CREAM if i % 2 == 0 else SOFT_BG
        add_rect(s, emu(0.5), y, emu(12.3), emu(0.55), bg)
        x = emu(0.5)
        for j, (w, val) in enumerate(zip(widths, row)):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            pad = emu(0.2) if j == 0 else 0
            add_textbox(s, x + pad, y, emu(w) - pad, emu(0.55), val, 13, j == 0, INK,
                        align, MSO_ANCHOR.MIDDLE)
            x += emu(w)
    add_footer(s, 9, TOTAL)
    return s


def slide_funnel_bases(prs):
    s = blank_slide(prs)
    add_header(s, "4. Показатели по базам", "Откуда пришли звонки — и что из этого вышло")
    rows = [
        ["Источник", "Прозвон", "Успешн.", "Сделка", "Запись", "Демо", "Сделка %"],
        ["База ЗУП", "91", "46", "10", "0", "3", "11,0%"],
        ["База КП", "124", "60", "17", "10", "5", "13,7%"],
        ["База канбана", "12", "12", "0", "0", "0", "0%"],
        ["База вебинаров", "298", "132", "13", "1", "1", "4,4%"],
        ["Заявка с сайта", "4", "4", "4", "0", "0", "100%"],
        ["База вахта", "125", "34", "1", "0", "0", "0,8%"],
        ["База общепит", "266", "124", "13", "0", "0", "4,9%"],
        ["Холодка общепит", "31", "12", "3", "0", "0", "9,7%"],
        ["База производства", "123", "54", "8", "1", "2", "6,5%"],
        ["Новая КС / прочие", "10", "5", "1", "0", "0", "10%"],
        ["Итого", "1 084", "483", "70", "12", "11", "6,5%"],
    ]
    col_w = [2.7, 1.5, 1.5, 1.5, 1.5, 1.5, 1.6]
    table_w = sum(col_w)
    start_x = emu(0.7)
    row_h = 0.38
    y0 = 1.82
    for r, row in enumerate(rows):
        y = emu(y0 + r * row_h)
        if r == 0:
            bg = INK
            color = WHITE
            bold = True
        elif r == len(rows) - 1:
            bg = YELLOW
            color = INK
            bold = True
        else:
            bg = CREAM if r % 2 else WHITE
            color = INK
            bold = False
        add_rect(s, start_x, y, emu(table_w), emu(row_h), bg)
        x = start_x
        for c, (w, val) in enumerate(zip(col_w, row)):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            pad = emu(0.12) if c == 0 else 0
            add_textbox(s, x + pad, y, emu(w) - pad, emu(row_h), val, 11, bold, color,
                        align, MSO_ANCHOR.MIDDLE)
            x += emu(w)
    add_textbox(
        s, emu(0.5), emu(6.55), emu(12.3), emu(0.45),
        "Узкое место воронки — не контакт, а запись и проведение демо. "
        "КП и ЗУП дают лучшую сделку; общепит и вебинары дают объём прозвона.",
        12, False, MUTED,
    )
    add_footer(s, 10, TOTAL)
    return s


def slide_cases(prs):
    s = blank_slide(prs)
    add_header(s, "5. Кейсы с демонстраций", "Сюда вписываем результаты живых демо — блок оставлен пустым")
    add_textbox(
        s, emu(0.5), emu(1.72), emu(12.3), emu(0.4),
        "За период проведено 11 демо. Ниже — шаблон карточки: клиент, что показали, чем закончилось.",
        13, False, MUTED,
    )
    fields = "Клиент / ИНН\nДата и менеджер\nЧто смотрели на демо\nРезультат и возражения\nСледующий шаг"
    for i in range(3):
        left = emu(0.5) + emu(i * 4.2)
        add_card(s, left, emu(2.25), emu(4.0), emu(4.55), CREAM, 0.08)
        add_rect(s, left, emu(2.25), emu(4.0), emu(0.55), YELLOW)
        add_textbox(s, left + emu(0.2), emu(2.32), emu(3.6), emu(0.42), f"Кейс {i+1}", 16, True, INK,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        add_textbox(s, left + emu(0.25), emu(2.95), emu(3.5), emu(3.55), fields, 13, False, RGBColor(0x9A, 0x9A, 0x9A))
    add_footer(s, 11, TOTAL)
    return s


def slide_next(prs):
    s = blank_slide(prs)
    add_header(s, "Дальше", "Что усиливаем после этого среза")
    steps = [
        ("1", "Дожать запись", "Сделки уже есть (70), записей 12. Скрипт H1 «входит в ИТС ПРОФ» — на каждую сделку своих."),
        ("2", "Пометить связку", "Выделить в аналитике источник «ДОКИ + КС», иначе H3 невозможно честно закрыть."),
        ("3", "Собрать кейсы демо", "11 встреч уже прошло — заполняем три карточки на предыдущем слайде."),
        ("4", "Сравнить контуры", "H4: вебинары vs первый обзвон. H2: целевой ОКВЭД vs нецелевой холод."),
    ]
    for i, (n, title, body) in enumerate(steps):
        y = emu(1.85) + emu(i * 1.15)
        add_card(s, emu(0.5), y, emu(12.3), emu(1.05), CREAM, 0.08)
        add_rect(s, emu(0.5), y, emu(1.05), emu(1.05), YELLOW)
        add_textbox(s, emu(0.5), y, emu(1.05), emu(1.05), n, 24, True, INK, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(s, emu(1.8), y + emu(0.15), emu(10.6), emu(0.35), title, 16, True, INK)
        add_textbox(s, emu(1.8), y + emu(0.52), emu(10.6), emu(0.42), body, 13, False, MUTED)
    add_footer(s, 12, TOTAL)
    return s


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_title(prs)
    slide_agenda(prs)
    slide_project(prs)
    slide_map(prs)
    slide_materials(prs)
    slide_hypotheses_overview(prs)
    slide_h_warm(prs)
    slide_h_cold(prs)
    slide_funnel_totals(prs)
    slide_funnel_bases(prs)
    slide_cases(prs)
    slide_next(prs)

    # blank layout exists on default pptx (layout 6). Confirm count.
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print("saved", OUT, "slides", len(prs.slides), "bytes", OUT.stat().st_size)


if __name__ == "__main__":
    main()
