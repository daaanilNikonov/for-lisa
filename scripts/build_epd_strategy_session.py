#!/usr/bin/env python3
"""Build yellow-white Forus-style presentation: strategic session on ЭПД."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "Стратегическая_сессия_ЭПД.pptx"
ASSETS = ROOT / "presentation" / "assets" / "epd_session"

# ГК Форус — жёлто-белый корпоративный стиль (шафран)
YELLOW = RGBColor(0xFE, 0xCF, 0x68)      # accent1 theme
YELLOW_DARK = RGBColor(0xF2, 0xAF, 0x2E)
YELLOW_SOFT = RGBColor(0xFF, 0xF4, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
GRAY_LIGHT = RGBColor(0x8A, 0x8A, 0x8A)
LINE = RGBColor(0xE8, 0xE8, 0xE8)
TEAL = RGBColor(0x39, 0x5F, 0x75)        # accent2 from Forus saffron theme

FONT = "Verdana"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def emu(inches: float) -> int:
    return int(Inches(inches))


def set_run(run, text, size_pt, bold=False, color=NEAR_BLACK, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }.get(anchor, "t"),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    size_pt=14,
    bold=False,
    color=NEAR_BLACK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        set_run(p.add_run(), text, size_pt, bold, color)
    return box


def add_paragraph(tf, text, size_pt=14, bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    set_run(p.add_run(), text, size_pt, bold, color)
    return p


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_round(slide, left, top, width, height, fill, corner=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_waves(slide, title=False):
    """Yellow corner ribbons in Forus promo style."""
    if title:
        path = ASSETS / "wave_title.png"
        if path.exists():
            slide.shapes.add_picture(str(path), emu(6.5), emu(2.8), emu(7.2), emu(4.8))
        path_tr = ASSETS / "wave_tr.png"
        if path_tr.exists():
            slide.shapes.add_picture(str(path_tr), emu(9.8), emu(-0.2), emu(3.8), emu(2.2))
        return

    tr = ASSETS / "wave_tr.png"
    bl = ASSETS / "wave_bl.png"
    if tr.exists():
        slide.shapes.add_picture(str(tr), emu(10.0), emu(-0.15), emu(3.6), emu(2.1))
    if bl.exists():
        slide.shapes.add_picture(str(bl), emu(-0.2), emu(5.5), emu(3.6), emu(2.1))


def add_footer(slide, page: int, total: int = 8):
    add_rect(slide, emu(0), emu(7.22), SLIDE_W, emu(0.28), YELLOW)
    add_textbox(
        slide,
        emu(0.6),
        emu(7.22),
        emu(8),
        emu(0.28),
        "ГК Форус  ·  Стратегическая сессия по ЭПД",
        size_pt=9,
        color=NEAR_BLACK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        emu(11.2),
        emu(7.22),
        emu(1.6),
        emu(0.28),
        f"{page} / {total}",
        size_pt=9,
        bold=True,
        color=NEAR_BLACK,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_brand_mark(slide, left=0.55, top=0.35):
    """Small yellow accent + FORUS wordmark."""
    add_round(slide, emu(left), emu(top), emu(0.28), emu(0.28), YELLOW, corner=0.5)
    add_textbox(
        slide,
        emu(left + 0.4),
        emu(top - 0.02),
        emu(2.5),
        emu(0.32),
        "ФОРУС",
        size_pt=12,
        bold=True,
        color=NEAR_BLACK,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_topic_badge(slide, number: int):
    badge = add_round(slide, emu(0.55), emu(1.05), emu(0.72), emu(0.72), YELLOW, corner=0.2)
    add_textbox(
        slide,
        emu(0.55),
        emu(1.05),
        emu(0.72),
        emu(0.72),
        str(number),
        size_pt=22,
        bold=True,
        color=NEAR_BLACK,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return badge


def add_discussion_card(slide, left, top, width, height, title, bullets):
    add_round(slide, left, top, width, height, YELLOW_SOFT, corner=0.08)
    # left yellow bar
    bar = add_rect(slide, left, top, emu(0.12), height, YELLOW)
    box = add_textbox(
        slide,
        left + emu(0.35),
        top + emu(0.22),
        width - emu(0.5),
        emu(0.4),
        title,
        size_pt=14,
        bold=True,
        color=NEAR_BLACK,
    )
    tf = box.text_frame
    for bullet in bullets:
        add_paragraph(tf, f"•  {bullet}", size_pt=13, color=GRAY, space_before=8)


def build_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)
    add_waves(slide, title=True)

    # Left yellow vertical accent
    add_rect(slide, emu(0), emu(0), emu(0.22), SLIDE_H, YELLOW)

    add_brand_mark(slide, 0.7, 0.55)

    add_textbox(
        slide,
        emu(0.7),
        emu(2.0),
        emu(10),
        emu(0.45),
        "СТРАТЕГИЧЕСКАЯ СЕССИЯ",
        size_pt=18,
        bold=True,
        color=TEAL,
    )
    add_textbox(
        slide,
        emu(0.7),
        emu(2.55),
        emu(11),
        emu(1.2),
        "ЭПД",
        size_pt=72,
        bold=True,
        color=NEAR_BLACK,
    )
    add_textbox(
        slide,
        emu(0.7),
        emu(3.85),
        emu(9),
        emu(0.7),
        "Вопросы к обсуждению: реклама, клиенты,\nобратная связь и финансовая модель направления",
        size_pt=18,
        color=GRAY,
    )

    # Yellow underline accent
    add_rect(slide, emu(0.7), emu(4.8), emu(2.4), emu(0.12), YELLOW)

    add_textbox(
        slide,
        emu(0.7),
        emu(5.2),
        emu(8),
        emu(0.4),
        "6 тем  ·  решения и следующие шаги",
        size_pt=14,
        color=GRAY_LIGHT,
    )
    add_footer(slide, 1)


def build_agenda(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_waves(slide)
    add_brand_mark(slide)
    add_textbox(
        slide,
        emu(0.55),
        emu(0.95),
        emu(10),
        emu(0.55),
        "Повестка сессии",
        size_pt=28,
        bold=True,
        color=NEAR_BLACK,
    )
    add_rect(slide, emu(0.55), emu(1.55), emu(1.6), emu(0.08), YELLOW)

    topics = [
        "Поле «время взятия в работу» заявки с РК",
        "Разделение стоимости РК между ЦП и ЦКС",
        "Продление РК и согласование бюджета",
        "ОС от МПП: проблемы, возражения, потребности",
        "Портрет клиента, который купил ЭПД",
        "Финансовая модель направления ЭПД",
    ]

    for i, topic in enumerate(topics):
        row = i // 2
        col = i % 2
        left = emu(0.55 + col * 6.2)
        top = emu(1.95 + row * 1.45)
        add_round(slide, left, top, emu(5.9), emu(1.2), YELLOW_SOFT, corner=0.1)
        add_round(slide, left + emu(0.25), top + emu(0.28), emu(0.55), emu(0.55), YELLOW, corner=0.2)
        add_textbox(
            slide,
            left + emu(0.25),
            top + emu(0.28),
            emu(0.55),
            emu(0.55),
            str(i + 1),
            size_pt=16,
            bold=True,
            color=NEAR_BLACK,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left + emu(1.0),
            top + emu(0.28),
            emu(4.6),
            emu(0.7),
            topic,
            size_pt=14,
            bold=True,
            color=NEAR_BLACK,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_footer(slide, 2)


def build_topic_slide(
    prs: Presentation,
    number: int,
    page: int,
    title: str,
    subtitle: str,
    blocks: list[tuple[str, list[str]]],
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_waves(slide)
    add_brand_mark(slide)
    add_topic_badge(slide, number)

    add_textbox(
        slide,
        emu(1.5),
        emu(1.1),
        emu(10.5),
        emu(0.7),
        title,
        size_pt=24,
        bold=True,
        color=NEAR_BLACK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_textbox(
            slide,
            emu(1.5),
            emu(1.8),
            emu(10.5),
            emu(0.4),
            subtitle,
            size_pt=13,
            color=TEAL,
        )

    n = len(blocks)
    gap = 0.25
    usable = 12.2
    card_w = (usable - gap * (n - 1)) / n
    top = emu(2.45)
    height = emu(4.2) if n <= 2 else emu(4.0)

    for i, (block_title, bullets) in enumerate(blocks):
        left = emu(0.55 + i * (card_w + gap))
        add_discussion_card(slide, left, top, emu(card_w), height, block_title, bullets)

    add_footer(slide, page)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title(prs)
    build_agenda(prs)

    build_topic_slide(
        prs,
        number=1,
        page=3,
        title="Поле «время взятия в работу» заявки с РК",
        subtitle="Операционный трекинг заявок из рекламной кампании",
        blocks=[
            (
                "Что обсуждаем",
                [
                    "Добавить поле фиксации времени\nвзятия заявки в работу",
                    "Источник заявок — рекламная\nкампания (РК)",
                    "Единый стандарт учёта для\nкоманды продаж / сопровождения",
                ],
            ),
            (
                "Вопросы к решению",
                [
                    "Где фиксируем: CRM / Битрикс / другое?",
                    "Кто заполняет и в какой момент?",
                    "Какие метрики считаем: SLA, скорость\nреакции, конверсия?",
                    "С какого числа включаем в процесс?",
                ],
            ),
        ],
    )

    build_topic_slide(
        prs,
        number=2,
        page=4,
        title="Разделение стоимости РК между ЦП и ЦКС",
        subtitle="Очная встреча 27.08 — согласование модели распределения затрат",
        blocks=[
            (
                "Контекст",
                [
                    "Рекламная кампания — общий\nинструмент для ЦП и ЦКС",
                    "Нужна прозрачная схема\nразделения бюджета",
                    "Решение фиксируем на очной\nвстрече 27 августа",
                ],
            ),
            (
                "К обсуждению",
                [
                    "Варианты: 50/50, по лидам,\nпо выручке, по заявкам",
                    "Что входит в стоимость РК?",
                    "Как учитываем в P&L\nнаправлений?",
                    "Кто владелец бюджета\nи кто согласует изменения?",
                ],
            ),
        ],
    )

    build_topic_slide(
        prs,
        number=3,
        page=5,
        title="Продление РК и согласование бюджета",
        subtitle="Текущая рекламная кампания действует до середины сентября",
        blocks=[
            (
                "Ситуация",
                [
                    "Текущая РК — до середины\nсентября",
                    "Нужно решение: продлеваем\nили останавливаем",
                    "Без согласования бюджета\nесть риск паузы в лидогенерации",
                ],
            ),
            (
                "Решения на сессии",
                [
                    "Продлеваем / не продлеваем?",
                    "На какой срок и с каким бюджетом?",
                    "Какие каналы оставляем / меняем?",
                    "Кто готовит и кто утверждает\nбюджет на следующий период?",
                ],
            ),
        ],
    )

    build_topic_slide(
        prs,
        number=4,
        page=6,
        title="ОС от МПП по сервису ЭПД",
        subtitle="Что говорят клиенты: проблемы, возражения и потребности",
        blocks=[
            (
                "Задача",
                [
                    "Собрать обратную связь от МПП",
                    "Зафиксировать типичные\nвозражения клиентов",
                    "Выделить основные проблемы\nи неудовлетворённые потребности",
                ],
            ),
            (
                "Формат сбора",
                [
                    "Короткий опрос / интервью МПП",
                    "Единый шаблон: проблема →\nвозражение → потребность",
                    "Срок сбора и ответственный",
                    "Как используем ОС: продукт,\nпродажи, маркетинг",
                ],
            ),
        ],
    )

    build_topic_slide(
        prs,
        number=5,
        page=7,
        title="Портрет клиента, который купил ЭПД",
        subtitle="Сегментация платящих клиентов для точного таргета и продаж",
        blocks=[
            (
                "Что собираем",
                [
                    "Отрасль, размер, география",
                    "Роль ЛПР и триггер покупки",
                    "Типичный сценарий сделки\nи средний чек",
                    "Почему выбрали именно ЭПД",
                ],
            ),
            (
                "Зачем это нужно",
                [
                    "Точнее настроить РК и офферы",
                    "Усилить скрипты МПП",
                    "Отделить «горячих» от\nнецелевых лидов",
                    "Связать портрет с финмоделью\nи прогнозом конверсий",
                ],
            ),
        ],
    )

    build_topic_slide(
        prs,
        number=6,
        page=8,
        title="Финансовая модель направления ЭПД",
        subtitle="Доходы, затраты, конверсии и сценарии — включая возможный перенос штрафов",
        blocks=[
            (
                "Доходы",
                [
                    "Титулы",
                    "Услуги",
                    "Продления",
                ],
            ),
            (
                "Затраты",
                [
                    "Реклама (РК)",
                    "Зарплаты",
                    "Производство / внедрение",
                ],
            ),
            (
                "Прогноз и сценарии",
                [
                    "Конверсии по воронке",
                    "Базовый / оптимистичный /\nпессимистичный",
                    "С учётом возможного\nпереноса штрафов",
                    "Срок подготовки модели\nи владелец",
                ],
            ),
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
