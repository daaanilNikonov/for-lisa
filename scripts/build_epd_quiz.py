#!/usr/bin/env python3
"""Build interactive «100 к 1 / Своя игра» quiz for 1С-ЭПД / Доки.Логистика."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Презентация ГК Форус темный шаблон 16х9 (1).pptx"
OUT_DIR = ROOT / "квиз 1с-эпд"
OUT = OUT_DIR / "Квиз_1С-ЭПД_Доки_Логистика_100к1.pptx"
OUT_COPY = ROOT / "presentation" / "quiz" / "Квиз_1С-ЭПД_Доки_Логистика_100к1.pptx"
TIMER_GIF = ROOT / "presentation" / "quiz" / "timer_30s.gif"

# Brand
BLUE = RGBColor(0x26, 0xA6, 0xE0)
CARD = RGBColor(0x3F, 0x3F, 0x3F)
CARD_DARK = RGBColor(0x2A, 0x2A, 0x2A)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xBF, 0xBF, 0xBF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GOLD = RGBColor(0xF0, 0xB4, 0x2E)
GREEN = RGBColor(0x3D, 0xC4, 0x8A)

FONT = "Verdana"
POINTS = [10, 30, 50, 70, 100]

L_TITLE = 0
L_CONTENT = 3
L_BG = 6

# ─── Questions: 5 topics × 5 points ─────────────────────────────────────────
# Source: «своя игра эпд.docx»

TOPICS: list[dict] = [
    {
        "name": "Основы ЭПД",
        "short": "Основы",
        "questions": [
            {
                "points": 10,
                "q": "Что означает сокращение ЭПД?",
                "options": [
                    "Электронные платёжные документы",
                    "Электронные перевозочные документы",
                    "Электронные первичные документы",
                ],
                "answer": "Электронные перевозочные документы",
            },
            {
                "points": 30,
                "q": "Что такое титул?",
                "options": [
                    "Отдельный документ для водителя в мобильном приложении",
                    "Часть ЭПД, которая заполняется конкретным участником перевозки в определённый момент",
                    "Электронная подпись участника перевозки",
                    "QR-код документа",
                ],
                "answer": "Часть ЭПД, которую заполняет конкретный участник перевозки в определённый момент",
            },
            {
                "points": 50,
                "q": "Какие участники подписывают ЭТрН?",
                "options": [
                    "Только грузоотправитель",
                    "Грузоотправитель, перевозчик и грузополучатель",
                    "Только перевозчик и водитель",
                    "Экспедитор и водитель",
                ],
                "answer": "Грузоотправитель, перевозчик и грузополучатель",
            },
            {
                "points": 70,
                "q": "Кто такой грузоотправитель?",
                "options": [
                    "Организация, которая принимает груз",
                    "Организация, которая владеет грузом и нанимает перевозчика",
                    "Организация, которая предоставляет транспорт",
                    "Человек, который совершает погрузку груза в ТС",
                ],
                "answer": "Организация, которая владеет грузом и нанимает перевозчика",
            },
            {
                "points": 100,
                "q": "Кто относится к основным участникам грузоперевозки?",
                "options": None,
                "answer": "Грузоотправитель, грузоперевозчик, грузополучатель, экспедитор",
            },
        ],
    },
    {
        "name": "ЭТрН и обязательность",
        "short": "ЭТрН",
        "questions": [
            {
                "points": 10,
                "q": "С какой даты для всех участников грузоперевозок становится обязательным использование документов в электронном виде?",
                "options": [
                    "1 января 2027 года",
                    "1 мая 2026 года",
                    "1 сентября 2026 года",
                    "1 марта 2027 года",
                ],
                "answer": "1 сентября 2026 года",
            },
            {
                "points": 30,
                "q": "В каком случае ЭТрН не требуется?",
                "options": [
                    "Все участники — юридические лица или ИП",
                    "Перевозчик оказывает коммерческую услугу",
                    "Хотя бы один участник — физическое лицо или самозанятый",
                    "Перевозчик нанят грузоотправителем",
                ],
                "answer": "Хотя бы один участник — физическое лицо или самозанятый",
            },
            {
                "points": 50,
                "q": "Компания доставляет собственным транспортом свой товар до покупателя. Нужно ли оформлять ЭТрН?",
                "options": [
                    "Да, всегда",
                    "Нет",
                    "Только если стоимость товара больше 100 000 рублей",
                    "Да, если перевозка по территории РФ",
                ],
                "answer": "Да, если перевозка по территории РФ",
            },
            {
                "points": 70,
                "q": "Компания продаёт товар, а покупатель сам забирает его своим транспортом. Нужна ли ЭТрН? Почему?",
                "options": None,
                "answer": "ЭТрН не требуется: перевозчик совпадает с грузополучателем (самовывоз)",
            },
            {
                "points": 100,
                "q": "Назовите условия, при которых ЭТрН обязательна к оформлению с 1 сентября 2026 года.",
                "options": None,
                "answer": (
                    "Все участники — ЮЛ или ИП; перевозчик не является грузоотправителем "
                    "и грузополучателем; перевозчик работает на коммерческой основе"
                ),
            },
        ],
    },
    {
        "name": "Документы и QR",
        "short": "Документы",
        "questions": [
            {
                "points": 10,
                "q": "Для чего водителю нужен QR-код?",
                "options": [
                    "Для оплаты перевозки",
                    "Для предъявления электронных перевозочных документов при проверке",
                    "Для получения путевого листа",
                ],
                "answer": "Для предъявления электронных перевозочных документов при проверке",
            },
            {
                "points": 30,
                "q": "Что такое ЭЗЗ и для чего она нужна?",
                "options": None,
                "answer": (
                    "Электронный заказ (заявка) — фиксирует договорённость между "
                    "грузоотправителем и перевозчиком до начала перевозки"
                ),
            },
            {
                "points": 50,
                "q": "Что содержит электронный путевой лист (ЭПЛ)?",
                "options": [
                    "Сведения о грузе, отправителе и адресах погрузки/выгрузки",
                    "Данные о водителе, ТС, маршруте, медосмотре, техконтроле и пробеге",
                    "Стоимость перевозки",
                    "Данные грузополучателя и грузоотправителя",
                ],
                "answer": "Данные о водителе, ТС, маршруте, медосмотре, техконтроле и пробеге",
            },
            {
                "points": 70,
                "q": "Нужен ли водителю интернет непосредственно во время проверки QR-кода на дороге?",
                "options": [
                    "Да, обязательно",
                    "Нет, интернет нужен для получения QR-кода, а во время проверки — нет",
                    "Интернет водителю не нужен вообще",
                ],
                "answer": "Нет, интернет нужен для получения QR-кода, а во время проверки — нет",
            },
            {
                "points": 100,
                "q": (
                    "Назовите хотя бы три сценария-исключения, при которых транспортная "
                    "накладная, заказ и заявка могут оформляться на бумаге."
                ),
                "options": None,
                "answer": (
                    "Военные перевозки; перевозки с иностранными юрлицами; каботаж ЕАЭС; "
                    "личные нужды физлица; сбой / нет интернета в утверждённых регионах; "
                    "отправитель = получатель; драгметаллы и камни и др."
                ),
            },
        ],
    },
    {
        "name": "Экспедиция и процесс",
        "short": "Процесс",
        "questions": [
            {
                "points": 10,
                "q": "В чём главное отличие экспедитора от агента в контексте ЭПД?",
                "options": [
                    "Экспедитор всегда является перевозчиком",
                    "Экспедитор принимает груз во владение и фигурирует в ЭТрН, агент — нет",
                    "Агент всегда подписывает экспедиторскую расписку",
                    "Отличий нет",
                ],
                "answer": "Экспедитор принимает груз во владение и фигурирует в ЭТрН, агент — нет",
            },
            {
                "points": 30,
                "q": (
                    "Поставщик обратился к экспедитору, экспедитор принял груз во владение "
                    "и нанял перевозчика. Какие роли возникают у участников?"
                ),
                "options": None,
                "answer": (
                    "Заказчик и экспедитор подписывают экспедиторские документы; "
                    "экспедитор принимает груз во владение, организует перевозку и нанимает "
                    "перевозчика; в ЭТрН экспедитор фигурирует как грузоотправитель"
                ),
            },
            {
                "points": 50,
                "q": (
                    "Цепочка: заказчик → экспедитор → перевозчик → получатель. "
                    "Какие электронные документы могут возникнуть у экспедитора?"
                ),
                "options": [
                    "Только ЭТрН и экспедиторская расписка",
                    "Экспедиторские документы, ЭЗЗ и ЭТрН",
                    "Только ЭПЛ",
                    "Только ЭЗЗ и ЭТрН",
                ],
                "answer": "Экспедиторские документы, ЭЗЗ и ЭТрН",
            },
            {
                "points": 70,
                "q": (
                    "Расставьте этапы перевозки в правильном порядке:\n"
                    "1) Приёмка груза получателем\n"
                    "2) Создание ЭТрН и отправка груза\n"
                    "3) Водитель в дороге\n"
                    "4) Создание и обработка ЭЗЗ\n"
                    "5) Разгрузка на складе получателя"
                ),
                "options": None,
                "answer": "4 → 2 → 3 → 1 → 5",
            },
            {
                "points": 100,
                "q": (
                    "Что произойдёт с необходимостью оформления ЭТрН, если грузоотправитель "
                    "является одновременно грузополучателем?"
                ),
                "options": [
                    "ЭТрН обязательна",
                    "ЭТрН не оформляется",
                    "Можно остаться на бумаге, а можно перейти на ЭТрН",
                ],
                "answer": "Можно остаться на бумаге, а можно перейти на ЭТрН",
            },
        ],
    },
    {
        "name": "1С-ЭПД и Доки",
        "short": "Решения",
        "questions": [
            {
                "points": 10,
                "q": "Что необходимо для подключения 1С-ЭПД? Назовите основные требования.",
                "options": None,
                "answer": (
                    "Программа 1С; УКЭП и МЧД для ответственных; КриптоПро CSP с лицензией; "
                    "для перевозчиков — подписание водителей через мобильное приложение "
                    "или с телефона"
                ),
            },
            {
                "points": 30,
                "q": "Назовите основные шаги подключения сервиса 1С-ЭПД.",
                "options": None,
                "answer": (
                    "Подключить 1С-ЭДО → обменяться приглашениями с контрагентами → "
                    "проверить настройку обмена ЭПД в учётной системе"
                ),
            },
            {
                "points": 50,
                "q": "Клиент говорит: «У нас вообще нет 1С, но нам нужно работать с ЭПД». Какое решение предложить?",
                "options": None,
                "answer": "1С:Клиент ЭДО или Доки",
            },
            {
                "points": 70,
                "q": "В чём ключевое отличие 1С-ЭПД и Доки с точки зрения сценариев работы клиента?",
                "options": None,
                "answer": (
                    "1С-ЭПД — типовое решение внутри конфигураций 1С. "
                    "Доки работает в 1С, веб и мобильном приложении независимо "
                    "и подходит клиентам без 1С, без ПК или с редко обновляемой 1С"
                ),
            },
            {
                "points": 100,
                "q": (
                    "Клиент: «С 1 сентября нам обязательно переходить на ЭПД. С чего начать?» "
                    "Дайте пошаговую рекомендацию."
                ),
                "options": None,
                "answer": (
                    "1) Проанализировать процессы  2) Выбрать ключевых контрагентов  "
                    "3) Подготовить ИС  4) Подключить сервис и сделать первую отправку  "
                    "5) Обучить персонал и обновить регламенты"
                ),
            },
        ],
    },
]


def emu(inches: float) -> int:
    return int(Inches(inches))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_run(run, text, size_pt, bold=False, color=WHITE, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color)
    return box


def add_card(slide, left, top, width, height, fill=CARD, corner=0.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def clear_body_placeholders(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (13, 14, 1, 2):
            if ph.has_text_frame:
                ph.text_frame.clear()


def fill_title(slide, text, size=26):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), text, size, True, WHITE)
            return ph
    return add_textbox(slide, emu(0.97), emu(0.48), emu(11.4), emu(1.0), text, size, True, WHITE)



def fill_shape_text(shape, text, size_pt=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    """Write text into a shape so the whole shape stays clickable."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size_pt, bold, color)
    return shape

def link_to_slide(shape, target_slide):
    """Create an internal hyperlink from shape to target_slide."""
    shape.click_action.target_slide = target_slide


RED = RGBColor(0xC0, 0x39, 0x2B)


def set_run_scheme_color(run, scheme: str = "hlink"):
    """Use theme hyperlink color so PowerPoint can switch it to folHlink after click."""
    r_pr = run._r.get_or_add_rPr()
    # Remove existing solidFill (any namespace form)
    for child in list(r_pr):
        if child.tag == qn("a:solidFill") or child.tag.endswith("}solidFill"):
            r_pr.remove(child)
    solid = etree.SubElement(r_pr, qn("a:solidFill"))
    scheme_el = etree.SubElement(solid, qn("a:schemeClr"))
    scheme_el.set("val", scheme)
    # Prevent python-pptx / Office from treating it as fixed sRGB
    try:
        run.font.color.theme_color = None
    except Exception:
        pass


def fill_shape_text_hyperlink(shape, text, size_pt=20, bold=True):
    """Score cell text bound to theme hyperlink colors (turns red when followed)."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = FONT
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", FONT)
    set_run_scheme_color(run, "hlink")
    # Second pass: font API sometimes rewrites solidFill to sRGB
    set_run_scheme_color(run, "hlink")
    return shape


def patch_theme_followed_hyperlink_red(prs: Presentation) -> None:
    """Set theme hyperlink / followed-hyperlink colors so opened cells turn red."""
    for part in prs.part.package.iter_parts():
        name = str(getattr(part, "partname", ""))
        if "theme" not in name or not name.endswith(".xml"):
            continue
        try:
            root = etree.fromstring(part.blob)
        except Exception:
            continue
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        changed = False
        hlink = root.find(".//a:clrScheme/a:hlink", ns)
        fol = root.find(".//a:clrScheme/a:folHlink", ns)
        if hlink is not None:
            for child in list(hlink):
                hlink.remove(child)
            srgb = etree.SubElement(hlink, qn("a:srgbClr"))
            srgb.set("val", "FFFFFF")  # unused cells: white score text
            changed = True
        if fol is not None:
            for child in list(fol):
                fol.remove(child)
            srgb = etree.SubElement(fol, qn("a:srgbClr"))
            srgb.set("val", "C0392B")  # opened cells: red
            changed = True
        if changed:
            part._blob = etree.tostring(
                root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )


def mark_cell_fill_red_on_follow(shape) -> None:
    """Also tint shape fill toward red via highlight — keep base fill, add red line."""
    try:
        shape.line.color.rgb = RED
        shape.line.width = Pt(1.5)
    except Exception:
        pass


def bind_text_run_hyperlink(shape) -> None:
    """Copy shape-level slide jump onto the text run (needed for followed-color)."""
    sp = shape._element
    hlink = None
    # Prefer cNvPr level (shape action)
    for el in sp.iter():
        if el.tag == qn("a:hlinkClick"):
            hlink = el
            break
    if hlink is None:
        return
    r_id = hlink.get(qn("r:id"))
    action = hlink.get("action")
    if not r_id:
        return
    for r in sp.iter(qn("a:r")):
        r_pr = r.find(qn("a:rPr"))
        if r_pr is None:
            r_pr = etree.Element(qn("a:rPr"))
            r.insert(0, r_pr)
        for old in list(r_pr):
            if old.tag == qn("a:hlinkClick") or old.tag.endswith("}hlinkClick"):
                r_pr.remove(old)
        # Re-assert scheme hyperlink color on the run
        for child in list(r_pr):
            if child.tag == qn("a:solidFill") or child.tag.endswith("}solidFill"):
                r_pr.remove(child)
        solid = etree.SubElement(r_pr, qn("a:solidFill"))
        scheme_el = etree.SubElement(solid, qn("a:schemeClr"))
        scheme_el.set("val", "hlink")
        hl = etree.SubElement(r_pr, qn("a:hlinkClick"))
        hl.set(qn("r:id"), r_id)
        if action:
            hl.set("action", action)




def fill_shape_answer(shape, answer_text: str):
    """Title + answer body inside one shape (for a single animation target)."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.TOP)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    set_run(p0.add_run(), "Ответ", 12, True, BLUE)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(6)
    p1.alignment = PP_ALIGN.LEFT
    set_run(p1.add_run(), answer_text, 13, False, NEAR_BLACK)
    # left margin via inset
    try:
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)
    except Exception:
        pass
    return shape


def add_appear_after_ms(slide, shape, delay_ms: int = 30000) -> None:
    """Entrance animation: shape is hidden until delay_ms, then appears (slideshow)."""
    spid = str(shape.shape_id)
    sld = slide._element

    # Drop previous timing block if present
    for old in list(sld.findall(qn("p:timing"))):
        sld.remove(old)

    # Minimal valid timing tree for auto Appear after delay
    timing_xml = f"""
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="{delay_ms}"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:par>
                                  <p:cTn id="5" presetID="1" presetClass="entr" presetSubtype="0"
                                         fill="hold" grpId="0" nodeType="withEffect">
                                    <p:stCondLst>
                                      <p:cond delay="0"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                      <p:set>
                                        <p:cBhvr>
                                          <p:cTn id="6" dur="1" fill="hold">
                                            <p:stCondLst>
                                              <p:cond delay="0"/>
                                            </p:stCondLst>
                                          </p:cTn>
                                          <p:tgtEl>
                                            <p:spTgt spid="{spid}"/>
                                          </p:tgtEl>
                                          <p:attrNameLst>
                                            <p:attrName>style.visibility</p:attrName>
                                          </p:attrNameLst>
                                        </p:cBhvr>
                                        <p:to>
                                          <p:strVal val="visible"/>
                                        </p:to>
                                      </p:set>
                                    </p:childTnLst>
                                  </p:cTn>
                                </p:par>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    """
    timing = etree.fromstring(timing_xml)
    sld.append(timing)



def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "Квиз", 22, False, SOFT)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(8)
            set_run(p2.add_run(), "1С-ЭПД / Доки.Логистика", 32, True, BLUE)
            p3 = tf.add_paragraph()
            p3.space_before = Pt(16)
            set_run(p3.add_run(), "Формат «100 к 1»  ·  Своя игра", 16, True, WHITE)
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "5 тем  ·  25 вопросов  ·  от 10 до 100 баллов\n"
                "Выберите ячейку на поле — вопрос, таймер 30 сек, затем ответ",
                13,
                False,
                SOFT,
            )
    return slide


def build_board(prs):
    """Jeopardy-style board: topics × points with hyperlinks."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Игровое поле  ·  100 к 1", 26)
    clear_body_placeholders(slide)
    add_textbox(
        slide,
        emu(0.97),
        emu(1.2),
        emu(11.4),
        emu(0.35),
        "Клик по баллам → вопрос  ·  открытые ячейки краснеют  ·  таймер 30 сек  ·  «К полю» — назад",
        12,
        False,
        SOFT,
    )

    # Layout: topic column + 5 point columns
    topic_w = emu(2.55)
    cell_w = emu(1.7)
    cell_h = emu(0.85)
    gap_x = emu(0.12)
    gap_y = emu(0.12)
    left0 = emu(0.7)
    top0 = emu(1.7)

    # Header row with point values
    theme_hdr = add_card(slide, left0, top0, topic_w, cell_h, CARD_DARK, 0.08)
    fill_shape_text(theme_hdr, "Тема", 14, True, SOFT)
    for ci, pts in enumerate(POINTS):
        left = left0 + topic_w + gap_x + ci * (cell_w + gap_x)
        hdr = add_card(slide, left, top0, cell_w, cell_h, BLUE, 0.1)
        fill_shape_text(hdr, str(pts), 18, True, WHITE)

    # Store cell shapes for later linking: (topic_idx, q_idx) -> shape
    cell_shapes: dict[tuple[int, int], object] = {}

    for ti, topic in enumerate(TOPICS):
        top = top0 + (ti + 1) * (cell_h + gap_y)
        # Topic label
        topic_card = add_card(slide, left0, top, topic_w, cell_h, CARD, 0.08)
        fill_shape_text(topic_card, topic["name"], 12, True, WHITE)
        for qi, q in enumerate(topic["questions"]):
            left = left0 + topic_w + gap_x + qi * (cell_w + gap_x)
            # Unused: blue / gold. After opening, score text turns red (followed hyperlink).
            fill = GOLD if q["points"] == 100 else BLUE
            shape = add_card(slide, left, top, cell_w, cell_h, fill, 0.12)
            fill_shape_text_hyperlink(shape, str(q["points"]), 22, True)
            cell_shapes[(ti, qi)] = shape

    return slide, cell_shapes


def build_question_slide(prs, topic, qdata, board_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, topic["name"], 22)
    clear_body_placeholders(slide)

    # Points badge (top-right)
    badge = add_card(slide, emu(10.55), emu(0.4), emu(1.15), emu(0.55), GOLD, 0.2)
    fill_shape_text(badge, f"{qdata['points']}", 18, True, NEAR_BLACK)

    # 30-second countdown timer (animated GIF, plays once in slideshow)
    if TIMER_GIF.exists():
        slide.shapes.add_picture(str(TIMER_GIF), emu(10.45), emu(1.1), width=emu(1.85))
        add_textbox(
            slide,
            emu(10.35),
            emu(2.95),
            emu(2.05),
            emu(0.35),
            "Таймер 30 сек",
            10,
            False,
            SOFT,
            PP_ALIGN.CENTER,
        )
    else:
        tcard = add_card(slide, emu(10.45), emu(1.1), emu(1.85), emu(1.85), CARD, 0.15)
        fill_shape_text(tcard, "30", 36, True, GOLD)

    # Question card (leave room for timer on the right)
    add_card(slide, emu(0.7), emu(1.35), emu(9.5), emu(2.2), CARD, 0.08)
    add_textbox(
        slide,
        emu(0.95),
        emu(1.5),
        emu(9.05),
        emu(1.9),
        qdata["q"],
        17 if len(qdata["q"]) < 160 else 14,
        True,
        WHITE,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    # Options or open marker
    y = emu(3.75)
    if qdata["options"]:
        for i, opt in enumerate(qdata["options"]):
            add_textbox(
                slide,
                emu(0.95),
                y + i * emu(0.36),
                emu(9.2),
                emu(0.34),
                f"{i + 1})  {opt}",
                12,
                False,
                SOFT,
            )
        ans_top = y + len(qdata["options"]) * emu(0.36) + emu(0.12)
    else:
        add_textbox(
            slide,
            emu(0.95),
            y,
            emu(9.2),
            emu(0.34),
            "Открытый вопрос — обсудите ответ командой",
            12,
            False,
            GOLD,
        )
        ans_top = y + emu(0.5)

    # Keep answer block from overlapping back button / bottom
    max_ans_bottom = emu(6.55)
    ans_h = max(emu(1.15), min(emu(1.55), max_ans_bottom - ans_top))

    # Answer card — appears automatically after 30 seconds
    answer_card = add_card(slide, emu(0.7), ans_top, emu(9.5), ans_h, CARD_LIGHT, 0.08)
    fill_shape_answer(answer_card, qdata["answer"])
    add_appear_after_ms(slide, answer_card, delay_ms=30000)

    # Hint under timer
    add_textbox(
        slide,
        emu(10.35),
        emu(3.3),
        emu(2.05),
        emu(0.7),
        "Ответ появится\nчерез 30 секунд",
        10,
        False,
        GOLD,
        PP_ALIGN.CENTER,
    )

    # Back button
    back = add_card(slide, emu(10.45), emu(5.35), emu(1.85), emu(1.15), BLUE, 0.12)
    fill_shape_text(back, "← К полю", 13, True, WHITE)
    link_to_slide(back, board_slide)

    return slide



def write_vba_macros(cell_shapes: dict) -> None:
    """Optional VBA: paint cell fill red and jump to question (for .pptm)."""
    lines = [
        "Attribute VB_Name = \"QuizMacros\"",
        "Option Explicit",
        "",
        "' Импорт: PowerPoint → Вид → Макросы → редактор VBA → File → Import",
        "' Затем сохраните презентацию как .pptm и назначьте макросы на ячейки",
        "' (или используйте pptx: там номера открытых вопросов краснеют без макросов).",
        "",
        "Private Sub PaintRed(oShp As Shape)",
        "    On Error Resume Next",
        "    oShp.Fill.Visible = msoTrue",
        "    oShp.Fill.Solid",
        "    oShp.Fill.ForeColor.RGB = RGB(192, 57, 43)  ' #C0392B",
        "    oShp.TextFrame.TextRange.Font.Color.RGB = RGB(255, 255, 255)",
        "End Sub",
        "",
        "Sub BackToBoard()",
        "    SlideShowWindows(1).View.GotoSlide 2",
        "End Sub",
        "",
    ]
    # Question slides start at index 3 (1-based in PPT)
    for (ti, qi), _cell in sorted(cell_shapes.items()):
        # slide number: title=1, board=2, questions start at 3 in row-major topic order
        q_index = ti * 5 + qi  # 0..24
        slide_num = 3 + q_index
        sub = f"OpenCell_{ti}_{qi}"
        lines += [
            f"Sub {sub}(oShp As Shape)",
            f"    PaintRed oShp",
            f"    SlideShowWindows(1).View.GotoSlide {slide_num}",
            "End Sub",
            "",
        ]
    out = OUT_DIR / "QuizMacros.bas"
    out.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"VBA macros: {out}")


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    if not TIMER_GIF.exists():
        raise SystemExit(f"Timer GIF not found: {TIMER_GIF}. Generate it first.")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    OUT_COPY.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    delete_all_slides(prs)

    build_title(prs)
    board_slide, cell_shapes = build_board(prs)

    # Build question slides and wire hyperlinks
    for ti, topic in enumerate(TOPICS):
        for qi, qdata in enumerate(topic["questions"]):
            q_slide = build_question_slide(prs, topic, qdata, board_slide)
            cell = cell_shapes.get((ti, qi))
            if cell is not None:
                link_to_slide(cell, q_slide)
                bind_text_run_hyperlink(cell)

    # Opened questions turn red via theme "followed hyperlink" color
    patch_theme_followed_hyperlink_red(prs)

    # Name cells for optional VBA (permanent red FILL)
    for (ti, qi), cell in cell_shapes.items():
        try:
            cell.name = f"Cell_{ti}_{qi}"
        except Exception:
            pass

    prs.save(str(OUT))
    shutil.copy2(OUT, OUT_COPY)

    # Write VBA helper for permanent red fill (optional .pptm)
    write_vba_macros(cell_shapes)

    n_q = sum(len(t["questions"]) for t in TOPICS)
    print(f"Saved: {OUT}")
    print(f"Copy:  {OUT_COPY}")
    print(f"Slides: {len(prs.slides)} (title + board + {n_q} questions)")
    print(f"Topics: {', '.join(t['name'] for t in TOPICS)}")
    print("Opened cells: score text turns red after return to the board (followed hyperlink).")


if __name__ == "__main__":
    main()
