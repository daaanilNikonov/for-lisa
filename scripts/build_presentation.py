#!/usr/bin/env python3
"""Build 7-slide pilot presentation using ГК Форус dark brand template."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Презентация ГК Форус темный шаблон 16х9 (1).pptx"
OUT = ROOT / "presentation" / "Продвижение_1С_Кабинет_сотрудника.pptx"
ICONS = ROOT / "presentation" / "assets" / "icons"

# Brand colors (ГК Форус — тёмный шаблон)
BLUE = RGBColor(0x26, 0xA6, 0xE0)
CARD = RGBColor(0x3F, 0x3F, 0x3F)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GRAY = RGBColor(0x76, 0x76, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xBF, 0xBF, 0xBF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
ROW_A = RGBColor(0x2A, 0x2A, 0x2A)
ROW_B = RGBColor(0x1F, 0x1F, 0x1F)

FONT = "Verdana"

L_TITLE = 0       # 1_Титульный слайд
L_CONTENT = 3     # 1_Слайд с текстом
L_CONTENT2 = 4    # 2_Слайд с текстом
L_BG = 6          # 1_Слайд с фоном
L_EMPTY = 22      # Пустой слайд


def emu(inches: float) -> int:
    return int(Inches(inches))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_run(run, text, size_pt, bold=False, color=WHITE, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    font_name=FONT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color, font_name)
    return box


def add_card(slide, left, top, width, height, fill=CARD, corner=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def add_icon(slide, name: str, left, top, width, height):
    path = ICONS / name
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width, height)


def clear_body_placeholders(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (13, 14, 1, 2):
            if ph.has_text_frame:
                ph.text_frame.clear()


def fill_title(slide, text, size=28):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), text, size, True, WHITE)
            return ph
    return add_textbox(slide, emu(0.97), emu(0.48), emu(11.4), emu(1.0), text, size, True, WHITE)


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "Продвижение сервиса", 26, True, WHITE)
            p2 = tf.add_paragraph()
            set_run(p2.add_run(), "«1С:Кабинет сотрудника»", 30, True, BLUE)
            p3 = tf.add_paragraph()
            p3.space_before = Pt(10)
            set_run(
                p3.add_run(),
                "Стратегия «Благодарность + Спецпредложение»\nдля топ-клиентов",
                16,
                False,
                WHITE,
            )
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "Персонализированный подход для увеличения продаж\n"
                "и укрепления лояльности среди ключевых клиентов в Иркутске",
                13,
                False,
                SOFT,
            )
    add_icon(slide, "icon_69.png", emu(10.55), emu(0.5), emu(1.8), emu(1.8))
    return slide


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Время — деньги, особенно для бухгалтера", 26)
    clear_body_placeholders(slide)

    points = [
        (
            "icon_58.png",
            "Рост нагрузки",
            "Кадровое делопроизводство требует всё больше времени. Бухгалтеры и кадровики "
            "перегружены рутиной: расчётные листки, отпуска, больничные.",
        ),
        (
            "icon_25.png",
            "Цифровизация неизбежна",
            "Рынок требует перехода на КЭДО, но многие компании пока не решаются на этот шаг.",
        ),
        (
            "icon_27.png",
            "Работа с возражениями",
            "Прямые продажи встречают сопротивление. Сначала показываем заботу — затем предлагаем решение.",
        ),
        (
            "icon_70.png",
            "Целевая аудитория",
            "100 самых активных и лояльных клиентов — те, кто уже доверяет нам и с большей вероятностью откликнется.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.7)
    card_w, card_h = emu(5.7), emu(2.2)
    gap_x, gap_y = emu(0.25), emu(0.2)

    for i, (icon_name, title, body) in enumerate(points):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_icon(slide, icon_name, left + emu(0.22), top + emu(0.55), emu(0.75), emu(0.75))
        add_textbox(
            slide, left + emu(1.15), top + emu(0.28), card_w - emu(1.4), emu(0.4),
            title, 15, True, BLUE,
        )
        add_textbox(
            slide, left + emu(1.15), top + emu(0.75), card_w - emu(1.4), emu(1.25),
            body, 12, False, WHITE,
        )
    return slide


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, "«Тёплый контакт» и выгодное предложение", 26)
    clear_body_placeholders(slide)

    steps = [
        (
            "1",
            "icon_22.png",
            "«Приятный сюрприз»",
            "Бухгалтеру (ЛПР) отправляется благодарственное письмо в красивой рамке "
            "за долголетнее сотрудничество.\n\n"
            "Цель: выделиться на фоне спама и вызвать положительные эмоции.",
        ),
        (
            "2",
            "icon_55.png",
            "«Личный звонок»",
            "Через 2–3 дня — персональный созвон менеджера: "
            "«Как вам подарок? Хотим предложить персональный тест-драйв…»\n\n"
            "Цель: перевести внимание с благодарности на бизнес-выгоду.",
        ),
        (
            "3",
            "icon_69.png",
            "«Специальное предложение»",
            "45 дней бесплатного доступа к «1С:Кабинет сотрудника» "
            "с полной настройкой и обучением.\n\n"
            "Цель: снять барьер входа и показать пользу на практике.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.8)
    card_w, card_h, gap = emu(3.85), emu(4.6), emu(0.25)

    for i, (num, icon_name, title, body) in enumerate(steps):
        left = left0 + i * (card_w + gap)
        highlight = i == 2
        fill = BLUE if highlight else CARD_LIGHT
        title_c = WHITE if highlight else NEAR_BLACK
        body_c = WHITE if highlight else RGBColor(0x33, 0x33, 0x33)
        num_c = WHITE if highlight else BLUE

        add_card(slide, left, top0, card_w, card_h, fill, 0.08)
        add_textbox(
            slide, left + emu(0.25), top0 + emu(0.18), emu(1.0), emu(0.65),
            num, 40, True, num_c,
        )
        add_icon(slide, icon_name, left + emu(2.65), top0 + emu(0.22), emu(0.9), emu(0.9))

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + emu(0.25),
            top0 + emu(1.05),
            emu(1.2),
            emu(0.05),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE if highlight else BLUE
        bar.line.fill.background()

        add_textbox(
            slide, left + emu(0.25), top0 + emu(1.25), card_w - emu(0.5), emu(0.65),
            title, 15, True, title_c,
        )
        add_textbox(
            slide, left + emu(0.25), top0 + emu(1.95), card_w - emu(0.5), emu(2.4),
            body, 12, False, body_c,
        )
    return slide


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT2])
    fill_title(slide, "Ожидаемые результаты и ключевые выгоды", 26)
    clear_body_placeholders(slide)

    left, top = emu(0.97), emu(1.65)
    w, h = emu(5.75), emu(4.8)

    # Client value
    add_card(slide, left, top, w, h, CARD_LIGHT, 0.06)
    add_textbox(slide, left + emu(0.3), top + emu(0.22), w - emu(0.6), emu(0.4),
                "Для клиента", 17, True, BLUE)
    add_textbox(slide, left + emu(0.3), top + emu(0.6), w - emu(0.6), emu(0.3),
                "Ценность для продажи", 11, False, GRAY)

    client_items = [
        ("Снижение нагрузки на персонал", "Автоматизация рутины, обработка заявок быстрее на 30%"),
        ("Экономия бюджета", "Сокращение расходов на бумагу, печать и хранение до 70%"),
        ("Юридическая значимость", "Расчётные листки по ст. 136 ТК РФ с подтверждением получения"),
        ("Современность", "Мобильный доступ для сотрудников и прозрачность процессов"),
    ]
    y = top + emu(1.1)
    for title, desc in client_items:
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + emu(0.3), y + emu(0.05), emu(0.28), emu(0.28)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        add_textbox(
            slide, left + emu(0.3), y + emu(0.02), emu(0.28), emu(0.28),
            "✓", 11, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(slide, left + emu(0.75), y, w - emu(1.1), emu(0.28), title, 13, True, NEAR_BLACK)
        add_textbox(
            slide, left + emu(0.75), y + emu(0.3), w - emu(1.1), emu(0.42),
            desc, 11, False, RGBColor(0x44, 0x44, 0x44),
        )
        y += emu(0.85)

    # Our KPIs
    left2 = emu(6.95)
    add_card(slide, left2, top, w, h, CARD, 0.06)
    add_textbox(slide, left2 + emu(0.3), top + emu(0.22), w - emu(0.6), emu(0.4),
                "Для нас", 17, True, BLUE)
    add_textbox(slide, left2 + emu(0.3), top + emu(0.6), w - emu(0.6), emu(0.3),
                "KPI проекта", 11, False, SOFT)

    kpi = [
        ("icon_66.png", "Конверсия 15–20%", "Перевести в активные продажи не менее 15–20% от 100 клиентов"),
        ("icon_22.png", "Укрепление лояльности", "Повышение NPS среди ключевых клиентов"),
        ("icon_67.png", "Имидж партнёра", "Статус надёжного партнёра, который заботится о бизнесе клиентов"),
        ("icon_59.png", "Рост продаж", "Новые подключения сервиса «1С:Кабинет сотрудника»"),
    ]
    y = top + emu(1.1)
    for icon_name, title, desc in kpi:
        add_icon(slide, icon_name, left2 + emu(0.25), y, emu(0.48), emu(0.48))
        add_textbox(slide, left2 + emu(0.9), y, w - emu(1.2), emu(0.28), title, 13, True, WHITE)
        add_textbox(slide, left2 + emu(0.9), y + emu(0.3), w - emu(1.2), emu(0.42), desc, 11, False, SOFT)
        y += emu(0.85)
    return slide


def style_table(table):
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            tc_pr = cell._tc.get_or_add_tcPr()
            for child in list(tc_pr):
                if any(x in child.tag for x in ("solidFill", "blipFill", "gradFill")):
                    tc_pr.remove(child)
            solid = etree.SubElement(tc_pr, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            if r == 0:
                srgb.set("val", "26A6E0")
                color, bold, size = WHITE, True, 11
            elif r == len(table.rows) - 1:
                srgb.set("val", "3F3F3F")
                color, bold, size = WHITE, True, 12
            else:
                srgb.set("val", "2A2A2A" if r % 2 else "1F1F1F")
                color, bold, size = WHITE, c == 0, 11

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                text = p.text
                # rebuild runs for reliable formatting
                for el in list(p._p):
                    if el.tag.endswith("}r"):
                        p._p.remove(el)
                set_run(p.add_run(), text, size, bold, color)


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Бюджет на привлечение 100 клиентов", 26)
    clear_body_placeholders(slide)

    rows_data = [
        ["Статья расходов", "Стоимость за ед., руб.", "Кол-во", "Итого, руб."],
        ["Печать благодарственного письма (А4, цветное)", "~ 100", "100 шт", "10 000"],
        ["Рамка для письма (А4, пластик/алюминий)", "~ 120–200 (опт)", "100 шт", "12 000"],
        ["Отправка (заказное письмо по г. Иркутску)", "~ 100–150", "100 шт", "15 000"],
        ["Итого бюджет", "", "", "37 000"],
    ]

    table_shape = slide.shapes.add_table(
        len(rows_data), 4, emu(0.97), emu(1.7), emu(8.55), emu(3.15)
    )
    table = table_shape.table
    table.columns[0].width = emu(4.35)
    table.columns[1].width = emu(1.7)
    table.columns[2].width = emu(1.2)
    table.columns[3].width = emu(1.3)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)

    add_card(slide, emu(9.8), emu(1.7), emu(2.6), emu(3.15), BLUE, 0.1)
    add_icon(slide, "icon_66.png", emu(10.45), emu(1.9), emu(1.3), emu(1.3))
    add_textbox(slide, emu(9.95), emu(3.25), emu(2.3), emu(0.35),
                "Итого", 13, False, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, emu(9.95), emu(3.55), emu(2.3), emu(0.6),
                "37 000 ₽", 26, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, emu(9.95), emu(4.2), emu(2.3), emu(0.45),
                "на 100 топ-клиентов", 11, False, WHITE, PP_ALIGN.CENTER)

    add_textbox(
        slide,
        emu(0.97),
        emu(5.15),
        emu(11.4),
        emu(1.0),
        "Примечание: цены ориентировочные (тендеры и розница в Иркутске). "
        "Окончательная стоимость — после уточнения у поставщиков. "
        "Заказное письмо обеспечивает отслеживание доставки и подтверждение вручения.",
        11,
        False,
        SOFT,
    )
    return slide


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Дорожная карта проекта", 28)
    clear_body_placeholders(slide)

    stages = [
        (
            "01",
            "Подготовка",
            "1 неделя",
            "• Сегментация базы: 100 активных клиентов\n"
            "• Дизайн и текст благодарности\n"
            "• Поиск поставщиков, закупка",
        ),
        (
            "02",
            "Отправка",
            "1–2 дня",
            "• Печать, рамки, подписание\n"
            "• Отправка заказными письмами\n"
            "• по г. Иркутску",
        ),
        (
            "03",
            "Обработка звонков",
            "2 недели",
            "• Скрипты для менеджеров\n"
            "• Обзвон после получения\n"
            "• Вебинары / демонстрации",
        ),
        (
            "04",
            "Анализ результатов",
            "1 неделя",
            "• Подсчёт конверсий\n"
            "• Сбор обратной связи\n"
            "• Отчёт о результатах",
        ),
    ]

    line_top = emu(1.85)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(1.35), line_top, emu(10.5), emu(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    left0, card_w, gap = emu(0.97), emu(2.9), emu(0.2)
    card_top = emu(2.25)

    for i, (num, title, timing, body) in enumerate(stages):
        left = left0 + i * (card_w + gap)
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + card_w // 2 - emu(0.16),
            line_top - emu(0.14),
            emu(0.32),
            emu(0.32),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = BLUE
        node.line.fill.background()

        add_card(slide, left, card_top, card_w, emu(3.55), CARD, 0.08)
        add_textbox(
            slide, left + emu(0.18), card_top + emu(0.15), card_w - emu(0.36), emu(0.5),
            num, 26, True, BLUE,
        )
        add_textbox(
            slide, left + emu(0.18), card_top + emu(0.7), card_w - emu(0.36), emu(0.45),
            title, 13, True, WHITE,
        )
        add_card(slide, left + emu(0.18), card_top + emu(1.2), emu(1.55), emu(0.32), BLUE, 0.5)
        add_textbox(
            slide, left + emu(0.18), card_top + emu(1.2), emu(1.55), emu(0.32),
            timing, 10, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left + emu(0.18), card_top + emu(1.7), card_w - emu(0.36), emu(1.65),
            body, 11, False, SOFT,
        )
    return slide


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Простое решение для больших продаж", 26)
    clear_body_placeholders(slide)

    theses = [
        (
            "Эмоции → доверие",
            "Дарим приятные эмоции — получаем доверие и шанс презентовать полезный продукт.",
        ),
        (
            "Бюджет как баннер",
            "37 000 ₽ — стоимость одного небольшого рекламного баннера, "
            "но с более высокой потенциальной конверсией.",
        ),
        (
            "Помогаем, а не продаём",
            "Решаем реальные проблемы клиентов (загруженность, бумаги) и укрепляем отношения.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.6)
    for i, (title, body) in enumerate(theses):
        top = top0 + i * emu(1.1)
        add_card(slide, left0, top, emu(8.3), emu(1.0), CARD, 0.08)
        mark = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left0 + emu(0.22), top + emu(0.32), emu(0.36), emu(0.36)
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = BLUE
        mark.line.fill.background()
        add_textbox(
            slide, left0 + emu(0.22), top + emu(0.32), emu(0.36), emu(0.36),
            str(i + 1), 12, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left0 + emu(0.8), top + emu(0.18), emu(7.2), emu(0.32),
            title, 14, True, BLUE,
        )
        add_textbox(
            slide, left0 + emu(0.8), top + emu(0.52), emu(7.2), emu(0.4),
            body, 12, False, WHITE,
        )

    add_card(slide, emu(9.5), emu(1.6), emu(3.0), emu(4.2), BLUE, 0.1)
    add_icon(slide, "icon_64.png", emu(10.25), emu(1.9), emu(1.5), emu(1.7))
    add_textbox(
        slide, emu(9.7), emu(3.7), emu(2.6), emu(0.85),
        "ПРОШУ\nОДОБРИТЬ", 18, True, WHITE, PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, emu(9.7), emu(4.65), emu(2.6), emu(0.85),
        "Пилотный проект\nбюджет 37 000 ₽", 13, False, WHITE, PP_ALIGN.CENTER,
    )
    return slide


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    if not ICONS.exists():
        raise SystemExit(f"Icons not found: {ICONS}. Prepare transparent icons first.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    delete_all_slides(prs)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
