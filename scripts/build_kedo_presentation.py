#!/usr/bin/env python3
"""Build pre-demo presentation for 1С:Кабинет сотрудника (ГК Форус)."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
MATERIALS = ROOT / "для презентаций по кабинету сотрудника"
TEMPLATE = MATERIALS / "Презентация ГК Форус темный шаблон 16х9 (1).pptx"
OUT_DIR = MATERIALS
OUT = OUT_DIR / "1С_Кабинет_сотрудника_Преддемонстрация.pptx"
OUT_COPY = ROOT / "presentation" / "1С_Кабинет_сотрудника_Преддемонстрация.pptx"
ICONS = ROOT / "presentation" / "assets" / "icons"
SCREENS = ROOT / "presentation" / "assets" / "screens"
FUNC = SCREENS / "func"

# Brand colors (ГК Форус — тёмный шаблон)
BLUE = RGBColor(0x26, 0xA6, 0xE0)
CARD = RGBColor(0x3F, 0x3F, 0x3F)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GRAY = RGBColor(0x76, 0x76, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xBF, 0xBF, 0xBF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_OK = RGBColor(0x3D, 0xC4, 0x8A)
WARN = RGBColor(0xF0, 0xA5, 0x00)
BAD = RGBColor(0xE0, 0x5A, 0x5A)

FONT = "Verdana"

L_TITLE = 0
L_CONTENT = 3
L_CONTENT2 = 4
L_BG = 6
L_EMPTY = 22


def emu(inches: float) -> int:
    return int(Inches(inches))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_run(run, text, size_pt, bold=False, color=WHITE, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    font_name=FONT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color, font_name)
    return box


def add_multiline(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size_pt=13,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    line_space=1.15,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        set_run(p.add_run(), line, size_pt, bold, color)
    return box


def add_card(slide, left, top, width, height, fill=CARD, corner=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def add_icon(slide, name: str, left, top, width, height):
    path = ICONS / name
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width, height)


def add_picture(slide, name: str, left, top, width=None, height=None):
    path = SCREENS / name
    if not path.exists():
        path = FUNC / name
    if not path.exists():
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def clear_body_placeholders(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (13, 14, 1, 2):
            if ph.has_text_frame:
                ph.text_frame.clear()


def fill_title(slide, text, size=26):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), text, size, True, WHITE)
            return ph
    return add_textbox(slide, emu(0.97), emu(0.48), emu(11.4), emu(1.0), text, size, True, WHITE)


def add_subtitle(slide, text, top=emu(1.25)):
    return add_textbox(slide, emu(0.97), top, emu(11.4), emu(0.45), text, 14, False, SOFT)


def add_check_dot(slide, left, top, size=0.28):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, emu(size), emu(size))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    add_textbox(
        slide,
        left,
        top,
        emu(size),
        emu(size),
        "✓",
        10 if size < 0.32 else 12,
        True,
        WHITE,
        PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_question_bar(slide, text, top=emu(6.35)):
    add_card(slide, emu(0.97), top, emu(11.4), emu(0.55), BLUE, 0.08)
    add_textbox(
        slide,
        emu(1.2),
        top + emu(0.08),
        emu(11.0),
        emu(0.4),
        "Вопрос: " + text,
        13,
        True,
        WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def style_table(table, header_blue=True):
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            tc_pr = cell._tc.get_or_add_tcPr()
            for child in list(tc_pr):
                if any(x in child.tag for x in ("solidFill", "blipFill", "gradFill")):
                    tc_pr.remove(child)
            solid = etree.SubElement(tc_pr, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            if r == 0:
                srgb.set("val", "26A6E0")
                color, bold, size = WHITE, True, 11
            else:
                srgb.set("val", "2A2A2A" if r % 2 else "1F1F1F")
                color, bold, size = WHITE, c == 0, 11
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                text = p.text
                for el in list(p._p):
                    if el.tag.endswith("}r"):
                        p._p.remove(el)
                set_run(p.add_run(), text, size, bold, color)


# ───────────────────────── slides ─────────────────────────


def build_slide_1(prs):
    """Титульный."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "1С:Кабинет сотрудника", 30, True, BLUE)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(12)
            set_run(
                p2.add_run(),
                "Кадровый электронный документооборот\nбез лишних программ и сложного внедрения",
                16,
                False,
                WHITE,
            )
            p3 = tf.add_paragraph()
            p3.space_before = Pt(18)
            set_run(p3.add_run(), "Преддемонстрация сервиса", 14, True, SOFT)
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "Все кадровые вопросы — в привычной системе 1С\n"
                "для сотрудников, руководителей, кадровиков и бухгалтерии",
                13,
                False,
                SOFT,
            )
    if (FUNC / "devices_laptop_phone.png").exists():
        add_picture(slide, "devices_laptop_phone.png", emu(7.6), emu(2.0), width=emu(5.0))
    elif (SCREENS / "devices_full.png").exists():
        add_picture(slide, "devices_full.png", emu(7.6), emu(2.0), width=emu(5.0))
    elif (SCREENS / "phone_brand.png").exists():
        add_picture(slide, "phone_brand.png", emu(10.2), emu(1.8), height=emu(4.8))
    return slide



def build_slide_forus(prs):
    """О группе компаний Форус — по данным forus.ru/about."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "О группе компаний «Форус»", 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Переводим бизнес в «цифру»  ·  www.forus.ru")

    stats = [
        ("1992", "год основания\nопыт в ИТ"),
        ("600+", "сотрудников\nв команде"),
        ("10 000+", "клиентов\nна сопровождении"),
        ("18", "регионов в России\nи офис в Монголии"),
    ]
    left0, top0 = emu(0.97), emu(1.75)
    card_w, card_h, gap = emu(2.8), emu(1.85), emu(0.2)
    for i, (num, label) in enumerate(stats):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top0, card_w, card_h, BLUE if i % 2 == 0 else CARD, 0.1)
        add_textbox(
            slide, left + emu(0.1), top0 + emu(0.25), card_w - emu(0.2), emu(0.55),
            num, 26, True, WHITE, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.15), top0 + emu(0.95), card_w - emu(0.3), emu(0.7),
            label, 12, False, WHITE, PP_ALIGN.CENTER,
        )

    pride = [
        ("ТОП", "крупнейших ИТ-компаний России"),
        ("1 место", "в сети «1С:БухОбслуживание» с 2017"),
        ("ТОП-5", "дистрибьюторов 1С в России"),
        ("Центр компетенции", "по кадровому ЭДО (1С)"),
    ]
    top2 = emu(3.85)
    for i, (title, body) in enumerate(pride):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top2, card_w, emu(1.55), CARD, 0.1)
        add_textbox(
            slide, left + emu(0.12), top2 + emu(0.2), card_w - emu(0.24), emu(0.45),
            title, 13, True, BLUE, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.12), top2 + emu(0.7), card_w - emu(0.24), emu(0.65),
            body, 11, False, WHITE, PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        emu(0.97),
        emu(5.6),
        emu(11.4),
        emu(1.0),
        "Внедряем ИТ-решения, цифровизируем процессы, помогаем компаниям развиваться.\n"
        "Направления: внедрение и сопровождение 1С, цифровизация, облачные сервисы, ЭДО и КЭДО, обучение.",
        12,
        False,
        SOFT,
    )
    return slide


def build_slide_2(prs):
    """Почему компании переходят на КЭДО?"""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Почему компании переходят на КЭДО?", 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Бумажный кадровый документооборот тормозит работу всей компании")

    cards = [
        ("icon_58.png", "Документы долго согласовываются", "Согласования растягиваются на дни, а не на минуты"),
        ("icon_25.png", "Бумага требует хранения и печати", "Печать, полки, курьеры и риск потери оригиналов"),
        ("icon_70.png", "Удалённые сотрудники едут ради подписи", "Дистанционные сотрудники тратят время на поездки"),
        ("icon_27.png", "Рутина кадровиков и бухгалтерии", "Часы ежедневно уходят на заявления, расчётные и архив"),
    ]
    left0, top0 = emu(0.97), emu(1.85)
    card_w, card_h = emu(5.7), emu(1.85)
    gap_x, gap_y = emu(0.25), emu(0.18)
    for i, (icon, title, body) in enumerate(cards):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_icon(slide, icon, left + emu(0.22), top + emu(0.5), emu(0.7), emu(0.7))
        add_textbox(slide, left + emu(1.1), top + emu(0.3), card_w - emu(1.35), emu(0.45), title, 14, True, BLUE)
        add_textbox(slide, left + emu(1.1), top + emu(0.85), card_w - emu(1.35), emu(0.7), body, 12, False, WHITE)

    add_card(slide, emu(0.97), emu(5.85), emu(11.4), emu(0.95), BLUE, 0.08)
    add_textbox(slide, emu(1.2), emu(5.95), emu(3.2), emu(0.75), "До 70%", 36, True, WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(
        slide,
        emu(4.3),
        emu(6.05),
        emu(7.8),
        emu(0.55),
        "экономии времени после перехода на КЭДО",
        18,
        True,
        WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


def build_slide_3(prs):
    """Что такое 1С:Кабинет сотрудника."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, "Что такое 1С:Кабинет сотрудника", 26)
    clear_body_placeholders(slide)
    add_subtitle(
        slide,
        "Мобильное взаимодействие сотрудников и бухгалтерии по кадровым вопросам",
    )

    steps = [
        ("Сотрудник", "Заявления, отпуска,\nсправки, подпись"),
        ("Руководитель", "Согласование\nсо смартфона"),
        ("1С:Кабинет\nсотрудника", "Единое цифровое\nпространство"),
        ("Кадровик", "Документы 1С\nбез повторного ввода"),
        ("Бухгалтер", "Расчётные листки\nс подтверждением"),
    ]
    left0, top0 = emu(0.7), emu(2.0)
    card_w, card_h = emu(2.15), emu(2.55)
    gap = emu(0.18)
    for i, (title, body) in enumerate(steps):
        left = left0 + i * (card_w + gap)
        fill = BLUE if i == 2 else CARD
        add_card(slide, left, top0, card_w, card_h, fill, 0.1)
        add_textbox(
            slide,
            left + emu(0.1),
            top0 + emu(0.35),
            card_w - emu(0.2),
            emu(0.9),
            title,
            13,
            True,
            WHITE,
            PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + emu(0.1),
            top0 + emu(1.4),
            card_w - emu(0.2),
            emu(0.9),
            body,
            11,
            False,
            WHITE if i == 2 else SOFT,
            PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            add_textbox(
                slide,
                left + card_w - emu(0.05),
                top0 + emu(1.0),
                emu(0.28),
                emu(0.4),
                "→",
                18,
                True,
                BLUE,
                PP_ALIGN.CENTER,
            )

    add_card(slide, emu(0.97), emu(5.0), emu(11.4), emu(1.5), CARD_LIGHT, 0.08)
    add_textbox(
        slide,
        emu(1.25),
        emu(5.25),
        emu(10.9),
        emu(0.4),
        "Не требуется отдельная HR-платформа",
        18,
        True,
        BLUE,
    )
    add_textbox(
        slide,
        emu(1.25),
        emu(5.7),
        emu(10.9),
        emu(0.55),
        "Работа продолжается в привычной 1С. Сотрудники получают расчётный лист, пишут заявления и заказывают справки с телефона или компьютера — без лишних обращений в отдел кадров.",
        12,
        False,
        NEAR_BLACK,
    )
    return slide


def build_slide_4(prs):
    """Что умеет сервис — плитки + скриншоты."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT2])
    fill_title(slide, "Что умеет сервис", 28)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Ключевые возможности 1С:Кабинет сотрудника")

    tiles = [
        "Заявления",
        "Отпуска",
        "Отсутствия",
        "Расчётные листки",
        "Справки",
        "Персональные данные",
        "Кадровые документы",
        "Согласование",
        "Электронная подпись",
        "Госключ",
    ]
    left0, top0 = emu(0.97), emu(1.8)
    tw, th = emu(2.35), emu(0.78)
    gap_x, gap_y = emu(0.15), emu(0.14)
    for i, title in enumerate(tiles):
        col, row = i % 3, i // 3
        if row > 3:
            break
        left = left0 + col * (tw + gap_x)
        top = top0 + row * (th + gap_y)
        add_card(slide, left, top, tw, th, CARD, 0.12)
        add_check_dot(slide, left + emu(0.15), top + emu(0.25), 0.28)
        add_textbox(
            slide,
            left + emu(0.55),
            top + emu(0.2),
            tw - emu(0.7),
            emu(0.4),
            title,
            12,
            True,
            WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Real screenshots of the mobile app
    if (FUNC / "mobile_home_salary.png").exists():
        add_picture(slide, "mobile_home_salary.png", emu(8.15), emu(1.65), height=emu(4.85))
    if (FUNC / "mobile_vacation_request.png").exists():
        add_picture(slide, "mobile_vacation_request.png", emu(10.55), emu(1.65), height=emu(4.85))
    return slide


def build_slide_5(prs):
    """Почему сотрудникам нравится сервис."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Почему сотрудникам нравится сервис", 24)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Все кадровые вопросы — со смартфона")

    if (FUNC / "mobile_home_salary.png").exists():
        add_picture(slide, "mobile_home_salary.png", emu(0.95), emu(1.65), height=emu(4.55))
    elif (SCREENS / "phone_mockup.png").exists():
        add_picture(slide, "phone_mockup.png", emu(0.85), emu(1.7), height=emu(4.5))

    items = [
        "Посмотреть расчётный листок",
        "Подать заявление",
        "Узнать остаток отпуска",
        "Получить справку",
        "Подписать документ",
        "Найти любой документ",
        "Выпустить подпись самостоятельно",
    ]
    left, top0 = emu(4.3), emu(1.75)
    for i, text in enumerate(items):
        top = top0 + i * emu(0.58)
        add_card(slide, left, top, emu(8.05), emu(0.52), CARD, 0.1)
        add_check_dot(slide, left + emu(0.15), top + emu(0.12), 0.28)
        add_textbox(slide, left + emu(0.55), top + emu(0.08), emu(7.3), emu(0.35), text, 13, False, WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_question_bar(slide, "Как сейчас сотрудники получают расчётные листки?")
    return slide


def build_role_slide(prs, title, subtitle, items, question=None, icon="icon_69.png"):
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, title, 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, subtitle)

    left0, top0 = emu(0.97), emu(1.8)
    card_w, card_h = emu(5.7), emu(1.25)
    gap_x, gap_y = emu(0.25), emu(0.18)
    for i, text in enumerate(items):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_check_dot(slide, left + emu(0.2), top + emu(0.45), 0.32)
        add_textbox(
            slide,
            left + emu(0.7),
            top + emu(0.25),
            card_w - emu(0.95),
            emu(0.75),
            text,
            13,
            False,
            WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    if question:
        add_question_bar(slide, question)
    return slide


def build_slide_6(prs):
    return build_role_slide(
        prs,
        "Что получает кадровик",
        "Меньше рутины — больше контроля",
        [
            "Заявления автоматически превращаются в документы 1С",
            "Данные не нужно вводить повторно",
            "Документы не теряются",
            "Все задачи собираются в одном месте",
            "Меньше ошибок при оформлении",
            "Электронный архив всегда доступен",
        ],
        "Сколько времени кадровая служба тратит на обработку заявлений?",
    )


def build_slide_7(prs):
    return build_role_slide(
        prs,
        "Что получает бухгалтер",
        "Автоматизация без лишней работы",
        [
            "Автоматическая отправка расчётных листков",
            "Подтверждение получения",
            "Меньше обращений сотрудников",
            "Меньше печати",
            "Документы формируются прямо из 1С",
            "Меньше ручного ввода данных",
        ],
    )


def build_slide_8(prs):
    return build_role_slide(
        prs,
        "Что получает руководитель",
        "Полный контроль процессов",
        [
            "Согласование документов с телефона",
            "Статусы документов",
            "Видно, кто подписал, а кто нет",
            "Можно назначить заместителя",
            "Быстрее согласование отпусков",
            "Прозрачность процессов",
        ],
        "Как сейчас руководители согласовывают отпуска?",
    )


def build_slide_9(prs):
    return build_role_slide(
        prs,
        "Что получает ИТ-специалист",
        "Простое сопровождение без новых систем",
        [
            "Встроено в 1С — без отдельной HR-платформы",
            "Не нужны доп. клиентские лицензии 1С",
            "Нет прямой нагрузки на базу 1С",
            "Нет прямого доступа сотрудников к ИБ 1С",
            "Облачная и локальная версия",
            "Хранение документов внутри 1С",
        ],
    )


def build_slide_10(prs):
    """Сравнение с другими решениями."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Почему именно 1С:Кабинет сотрудника", 24)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Сравнение с другими решениями")

    rows = [
        ["Возможность", "1С КС", "HRlink", "VK HR Tek", "Контур КЭДО", "Saby"],
        ["Работа внутри 1С", "✅", "❌", "❌", "❌", "❌"],
        ["Мобильное приложение", "✅", "✅", "✅", "✅", "✅"],
        ["Электронные подписи", "✅", "✅", "✅", "✅", "✅"],
        ["Автосоздание документов", "✅", "⚠", "⚠", "⚠", "⚠"],
        ["Без двойной работы", "✅", "⚠", "⚠", "⚠", "⚠"],
    ]
    table_shape = slide.shapes.add_table(len(rows), 6, emu(0.7), emu(1.85), emu(12.0), emu(3.6))
    table = table_shape.table
    widths = [3.2, 1.55, 1.55, 1.7, 2.0, 1.5]
    for i, w in enumerate(widths):
        table.columns[i].width = emu(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)

    add_card(slide, emu(0.97), emu(5.7), emu(11.4), emu(1.1), CARD_LIGHT, 0.08)
    add_textbox(
        slide,
        emu(1.2),
        emu(5.85),
        emu(10.9),
        emu(0.35),
        "Главное преимущество — сотрудники продолжают работать в привычной системе 1С.",
        14,
        True,
        BLUE,
    )
    add_textbox(
        slide,
        emu(1.2),
        emu(6.25),
        emu(10.9),
        emu(0.4),
        "Сравнение ориентировочное и требует проверки перед публикацией.",
        11,
        False,
        GRAY,
    )
    return slide


def build_slide_11(prs):
    """Экономический эффект."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, "Экономический эффект", 28)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Реальная экономия по данным проектов и материалов Форус")

    stats = [
        ("До 70%", "экономии времени\nна кадровой рутине"),
        ("До 475 тыс. ₽", "экономии за три месяца\nна одном из проектов"),
        ("От 21 ₽", "за сотрудника\nв месяц"),
        ("Более 5000", "клиентов уже\nиспользуют сервис"),
    ]
    left0, top0 = emu(0.85), emu(1.95)
    card_w, card_h = emu(2.9), emu(3.5)
    gap = emu(0.2)
    for i, (num, label) in enumerate(stats):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top0, card_w, card_h, CARD if i % 2 == 0 else BLUE, 0.1)
        add_textbox(
            slide,
            left + emu(0.15),
            top0 + emu(0.7),
            card_w - emu(0.3),
            emu(1.1),
            num,
            22 if len(num) > 8 else 26,
            True,
            WHITE,
            PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + emu(0.2),
            top0 + emu(2.0),
            card_w - emu(0.4),
            emu(1.1),
            label,
            13,
            False,
            WHITE,
            PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        emu(0.97),
        emu(5.55),
        emu(11.4),
        emu(1.15),
        "Стоимость — от 21 до 30 ₽/мес. на сотрудника (forus.ru). "
        "10 кабинетов — 280 ₽/мес. · 25 — 700 ₽ · 50 — 1 400 ₽ · 100 — 2 800 ₽ при оплате за год.\n"
        "Кейсы: «Фанат Байкала» полностью ушёл от бумаги; DNS переводит 40 000 сотрудников на КЭДО.",
        12,
        False,
        SOFT,
    )
    return slide


def build_slide_12(prs):
    """Как проходит внедрение."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Как проходит внедрение", 28)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Простой путь к запуску без ломки процессов 1С")

    steps = [
        ("01", "Подключение", "Подключаем сервис\nк вашей 1С"),
        ("02", "Настройка", "Роли, процессы,\nпечатные формы"),
        ("03", "Выпуск подписей", "УНЭП / Госключ\nдля сотрудников"),
        ("04", "Обучение", "Видеоуроки\nи короткие инструкции"),
        ("05", "Начало работы", "Заявления, расчётные,\nсогласования"),
    ]
    # timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(1.4), emu(2.55), emu(10.5), emu(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    left0, card_w, gap = emu(0.7), emu(2.3), emu(0.18)
    for i, (num, title, body) in enumerate(steps):
        left = left0 + i * (card_w + gap)
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + card_w // 2 - emu(0.16),
            emu(2.4),
            emu(0.32),
            emu(0.32),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = BLUE
        node.line.fill.background()

        add_card(slide, left, emu(3.1), card_w, emu(2.4), CARD, 0.1)
        add_textbox(slide, left + emu(0.15), emu(3.25), card_w - emu(0.3), emu(0.4), num, 18, True, BLUE, PP_ALIGN.CENTER)
        add_textbox(slide, left + emu(0.15), emu(3.7), card_w - emu(0.3), emu(0.45), title, 13, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, left + emu(0.15), emu(4.3), card_w - emu(0.3), emu(0.95), body, 11, False, SOFT, PP_ALIGN.CENTER)

    add_card(slide, emu(0.97), emu(5.8), emu(11.4), emu(0.85), CARD_LIGHT, 0.08)
    add_textbox(
        slide,
        emu(1.2),
        emu(5.95),
        emu(10.9),
        emu(0.55),
        "Внедрение не меняет привычный процесс работы в 1С.",
        16,
        True,
        BLUE,
        PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide



def build_slide_screens_employee(prs):
    """Скриншоты: сервис глазами сотрудника — для быстрого показа менеджером."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Сервис глазами сотрудника", 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Покажите клиенту экраны мобильного приложения до демо")

    shots = [
        ("mobile_home_salary.png", "Главный экран:\nзарплата и отсутствия"),
        ("mobile_vacation_request.png", "Заявление\nна отпуск"),
        ("mobile_vacation_wishes.png", "Пожелания\nк графику отпусков"),
        ("mobile_manager.png", "Раздел\nруководителя"),
    ]
    left0, top0 = emu(0.55), emu(1.7)
    card_w, gap = emu(3.0), emu(0.2)
    for i, (img, caption) in enumerate(shots):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top0, card_w, emu(5.0), CARD, 0.08)
        path = FUNC / img
        if path.exists():
            # Fit image inside card
            add_picture(slide, img, left + emu(0.25), top0 + emu(0.2), height=emu(3.85))
        add_textbox(
            slide,
            left + emu(0.1),
            top0 + emu(4.15),
            card_w - emu(0.2),
            emu(0.7),
            caption,
            11,
            True,
            WHITE,
            PP_ALIGN.CENTER,
        )
    return slide


def build_slide_screens_hr(prs):
    """Скриншоты: работа кадровика в 1С."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Работа кадровика прямо в 1С", 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, "Заявления превращаются в задачи и документы без двойного ввода")

    # Top row: 2 wider 1C screenshots
    top_shots = [
        ("1c_vacation_task.png", "Задача по заявлению на отпуск"),
        ("1c_transfer_to_cabinet.png", "Кнопка «Передать в 1С:Кабинет сотрудника»"),
    ]
    left0, top0 = emu(0.7), emu(1.7)
    card_w, card_h, gap = emu(5.9), emu(2.55), emu(0.25)
    for i, (img, caption) in enumerate(top_shots):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top0, card_w, card_h, CARD, 0.08)
        if (FUNC / img).exists():
            add_picture(slide, img, left + emu(0.15), top0 + emu(0.12), width=emu(5.6))
        add_textbox(
            slide,
            left + emu(0.15),
            top0 + emu(2.1),
            card_w - emu(0.3),
            emu(0.35),
            caption,
            12,
            True,
            BLUE,
        )

    # Bottom row
    bot_shots = [
        ("1c_kedo_documents.png", "Журнал документов кадрового ЭДО"),
        ("1c_hire_print_forms.png", "Печатные формы при приёме → в КС"),
    ]
    top1 = emu(4.45)
    for i, (img, caption) in enumerate(bot_shots):
        left = left0 + i * (card_w + gap)
        add_card(slide, left, top1, card_w, emu(2.2), CARD, 0.08)
        if (FUNC / img).exists():
            add_picture(slide, img, left + emu(0.15), top1 + emu(0.12), width=emu(5.6))
        add_textbox(
            slide,
            left + emu(0.15),
            top1 + emu(1.75),
            card_w - emu(0.3),
            emu(0.35),
            caption,
            12,
            True,
            BLUE,
        )
    return slide


def build_slide_13(prs):
    """Что покажем на демонстрации."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Что покажем на демонстрации", 26)
    clear_body_placeholders(slide)
    add_subtitle(slide, "После экранов выше — живая демонстрация на ваших процессах")

    items = [
        ("icon_55.png", "Оформление отпуска"),
        ("icon_22.png", "Отправка расчётного листка"),
        ("icon_69.png", "Согласование руководителем"),
        ("icon_66.png", "Подписание документов"),
        ("icon_59.png", "Выпуск подписи"),
        ("icon_70.png", "Работа мобильного приложения"),
        ("icon_27.png", "Работа кадровика в 1С"),
    ]
    left0, top0 = emu(0.97), emu(1.85)
    card_w, card_h = emu(5.7), emu(1.0)
    gap_x, gap_y = emu(0.25), emu(0.15)
    for i, (icon, text) in enumerate(items):
        col, row = i % 2, i // 2
        if i == 6:
            left = emu(0.97) + emu(2.975)
            top = top0 + 3 * (card_h + gap_y)
        else:
            left = left0 + col * (card_w + gap_x)
            top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w if i < 6 else emu(5.7), card_h, CARD, 0.1)
        add_icon(slide, icon, left + emu(0.2), top + emu(0.2), emu(0.55), emu(0.55))
        add_textbox(
            slide,
            left + emu(0.95),
            top + emu(0.25),
            emu(4.5),
            emu(0.5),
            text,
            14,
            True,
            WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    return slide


def build_slide_14(prs):
    """Обсудим ваши процессы — открытый вопрос вместо «Спасибо»."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "Обсудим ваши процессы", 28, True, WHITE)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(18)
            set_run(
                p2.add_run(),
                "Какие кадровые задачи сегодня\nзанимают больше всего времени?",
                20,
                True,
                BLUE,
            )
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "Открытый вопрос помогает перейти к демонстрации\n"
                "исходя из реальных потребностей вашей компании.\n\n"
                "8 (3952) 78-00-00  ·  www.forus.ru/produkty-1c/servisy/1c-kabinet-sotrudnika/",
                13,
                False,
                SOFT,
            )
    add_icon(slide, "icon_64.png", emu(10.4), emu(0.55), emu(1.9), emu(1.9))
    return slide


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    if not ICONS.exists():
        raise SystemExit(f"Icons not found: {ICONS}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    delete_all_slides(prs)

    builders = [
        build_slide_1,
        build_slide_forus,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
        build_slide_11,
        build_slide_12,
        build_slide_screens_employee,
        build_slide_screens_hr,
        build_slide_13,
        build_slide_14,
    ]
    for build in builders:
        build(prs)

    prs.save(str(OUT))
    OUT_COPY.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(OUT, OUT_COPY)
    print(f"Saved: {OUT}")
    print(f"Copy:  {OUT_COPY}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
