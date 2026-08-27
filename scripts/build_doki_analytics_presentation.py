#!/usr/bin/env python3
"""Build analytical presentation: project «Доки» — defense & handoff for scale."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Презентация ГК Форус темный шаблон 16х9 (1).pptx"
OUT = ROOT / "presentation" / "Доки_Аналитика_защита_и_тиражирование.pptx"
ICONS = ROOT / "presentation" / "assets" / "icons"

# Brand colors (ГК Форус — тёмный шаблон)
BLUE = RGBColor(0x26, 0xA6, 0xE0)
CARD = RGBColor(0x3F, 0x3F, 0x3F)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GRAY = RGBColor(0x76, 0x76, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xBF, 0xBF, 0xBF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)

FONT = "Verdana"

L_TITLE = 0
L_CONTENT = 3
L_CONTENT2 = 4
L_BG = 6
L_EMPTY = 22
L_LAST = 23


def emu(inches: float) -> int:
    return int(Inches(inches))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_run(run, text, size_pt, bold=False, color=WHITE, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = etree.SubElement(r_pr, qn(tag))
        el.set("typeface", font_name)


def set_anchor(text_frame, anchor=MSO_ANCHOR.TOP):
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    font_name=FONT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    tf.paragraphs[0].alignment = align
    set_run(tf.paragraphs[0].add_run(), text, size_pt, bold, color, font_name)
    return box


def add_multiline(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size_pt=13,
    bold=False,
    color=WHITE,
    line_space=Pt(6),
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.TOP)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = line_space
        set_run(p.add_run(), line, size_pt, bold, color)
    return box


def add_card(slide, left, top, width, height, fill=CARD, corner=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    return shape


def add_icon(slide, name: str, left, top, width, height):
    path = ICONS / name
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width, height)


def clear_body_placeholders(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (13, 14, 1, 2):
            if ph.has_text_frame:
                ph.text_frame.clear()


def fill_title(slide, text, size=26):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), text, size, True, WHITE)
            return ph
    return add_textbox(slide, emu(0.97), emu(0.48), emu(11.4), emu(1.0), text, size, True, WHITE)


def style_table(table, emphasize_last=False):
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            tc_pr = cell._tc.get_or_add_tcPr()
            for child in list(tc_pr):
                if any(x in child.tag for x in ("solidFill", "blipFill", "gradFill")):
                    tc_pr.remove(child)
            solid = etree.SubElement(tc_pr, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            if r == 0:
                srgb.set("val", "26A6E0")
                color, bold, size = WHITE, True, 12
            elif emphasize_last and r == len(table.rows) - 1:
                srgb.set("val", "3F3F3F")
                color, bold, size = WHITE, True, 12
            else:
                srgb.set("val", "2A2A2A" if r % 2 else "1F1F1F")
                color, bold, size = WHITE, c == 0, 12

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                text = p.text
                for el in list(p._p):
                    if el.tag.endswith("}r"):
                        p._p.remove(el)
                set_run(p.add_run(), text, size, bold, color)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(p.add_run(), "Проект «Доки»", 28, True, BLUE)
            p2 = tf.add_paragraph()
            set_run(p2.add_run(), "Итоги 7 недель и передача в тиражирование", 24, True, WHITE)
            p3 = tf.add_paragraph()
            p3.space_before = Pt(12)
            set_run(
                p3.add_run(),
                "Аналитика для защиты проекта продаж\n"
                "и передачи другой группе",
                15,
                False,
                SOFT,
            )
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            set_run(
                p.add_run(),
                "2 гипотезы  ·  18 подключений  ·  оффер: 3 месяца безлимит бесплатно",
                13,
                False,
                SOFT,
            )
    add_icon(slide, "icon_66.png", emu(10.55), emu(0.5), emu(1.8), emu(1.8))
    return slide


def build_slide_2(prs):
    """Зачем делали проект."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Зачем запускали проект", 26)
    clear_body_placeholders(slide)

    cards = [
        (
            "icon_25.png",
            "Задача",
            "Проверить, можно ли стабильно продавать «Доки» "
            "двумя понятными способами — и собрать пакет для тиражирования.",
        ),
        (
            "icon_58.png",
            "Срок",
            "7 недель пилота: гипотезы, скрипты, материалы "
            "для рассылки и сравнения, подключения клиентов.",
        ),
        (
            "icon_70.png",
            "Результат для защиты",
            "Показать цифры, что сработало, что готово к передаче "
            "и какие вопросы ещё нужно закрыть перед масштабом.",
        ),
        (
            "icon_27.png",
            "Для кого сейчас",
            "Защита перед продажами и передача другой группе, "
            "которая будет тиражировать подход.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.65)
    card_w, card_h = emu(5.7), emu(2.25)
    gap_x, gap_y = emu(0.25), emu(0.2)

    for i, (icon_name, title, body) in enumerate(cards):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_icon(slide, icon_name, left + emu(0.22), top + emu(0.55), emu(0.7), emu(0.7))
        add_textbox(
            slide, left + emu(1.1), top + emu(0.28), card_w - emu(1.35), emu(0.4),
            title, 15, True, BLUE,
        )
        add_textbox(
            slide, left + emu(1.1), top + emu(0.75), card_w - emu(1.35), emu(1.3),
            body, 13, False, WHITE,
        )
    return slide


def build_slide_3(prs):
    """Ключевые цифры."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, "Ключевые цифры за 7 недель", 26)
    clear_body_placeholders(slide)

    metrics = [
        ("7", "недель", "длительность пилота"),
        ("2", "гипотезы", "проверены в поле"),
        ("18", "подключений", "общее число продаж"),
        ("9 / 9", "поровну", "обе гипотезы дали результат"),
    ]

    left0, top0 = emu(0.97), emu(1.75)
    card_w, card_h, gap = emu(2.9), emu(3.55), emu(0.2)

    for i, (num, label, note) in enumerate(metrics):
        left = left0 + i * (card_w + gap)
        highlight = i == 2
        fill = CARD_LIGHT if highlight else CARD
        num_c = BLUE if highlight else WHITE
        label_c = NEAR_BLACK if highlight else WHITE
        note_c = GRAY if highlight else SOFT
        add_card(slide, left, top0, card_w, card_h, fill, 0.1)
        add_textbox(
            slide, left + emu(0.15), top0 + emu(0.7), card_w - emu(0.3), emu(1.0),
            num, 36 if i < 3 else 28, True, num_c, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.15), top0 + emu(1.85), card_w - emu(0.3), emu(0.45),
            label, 16, True, label_c, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.2), top0 + emu(2.45), card_w - emu(0.4), emu(0.7),
            note, 12, False, note_c, PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        emu(0.97),
        emu(5.5),
        emu(11.4),
        emu(0.6),
        "Оффер на входе: бесплатное использование безлимитным пакетом на 3 месяца.",
        14,
        False,
        SOFT,
    )
    return slide


def build_slide_4(prs):
    """Гипотеза 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Гипотеза 1. Доки вместо 1С-ЭПД", 26)
    clear_body_placeholders(slide)

    add_card(slide, emu(0.97), emu(1.55), emu(11.4), emu(1.15), CARD, 0.08)
    add_textbox(
        slide, emu(1.2), emu(1.7), emu(11.0), emu(0.35),
        "Идея", 13, True, BLUE,
    )
    add_textbox(
        slide, emu(1.2), emu(2.05), emu(11.0), emu(0.5),
        "Подключать «Доки» клиентам из базы ЭПД, которым не подходит решение 1С-ЭПД.",
        15, False, WHITE,
    )

    left_items = [
        ("Кто клиент", "Компании, которым нужен ЭДО, но 1С-ЭПД слишком узкий или неудобный."),
        ("Почему Доки", "Универсальный обмен документами, веб и 1С, облачный архив."),
        ("Что продаём", "Альтернативу там, где логистический ЭПД не закрывает задачу."),
    ]

    y = emu(2.95)
    for title, body in left_items:
        add_card(slide, emu(0.97), y, emu(7.4), emu(0.85), CARD, 0.06)
        add_textbox(slide, emu(1.2), y + emu(0.12), emu(7.0), emu(0.28), title, 13, True, BLUE)
        add_textbox(slide, emu(1.2), y + emu(0.42), emu(7.0), emu(0.35), body, 12, False, WHITE)
        y += emu(0.95)

    add_card(slide, emu(8.6), emu(2.95), emu(3.75), emu(2.85), BLUE, 0.1)
    add_textbox(
        slide, emu(8.85), emu(3.2), emu(3.25), emu(0.35),
        "Итог гипотезы 1", 13, False, WHITE, PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, emu(8.85), emu(3.65), emu(3.25), emu(0.7),
        "9", 44, True, WHITE, PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, emu(8.85), emu(4.4), emu(3.25), emu(0.4),
        "продаж", 16, True, WHITE, PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, emu(8.95), emu(4.95), emu(3.05), emu(0.6),
        "Подтверждена:\nканал из базы ЭПД работает",
        12, False, WHITE, PP_ALIGN.CENTER,
    )
    return slide


def build_slide_5(prs):
    """Гипотеза 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Гипотеза 2. Доки в скрипте с другими сервисами", 24)
    clear_body_placeholders(slide)

    add_card(slide, emu(0.97), emu(1.5), emu(11.4), emu(1.05), CARD, 0.08)
    add_textbox(
        slide, emu(1.2), emu(1.62), emu(11.0), emu(0.28),
        "Идея", 13, True, BLUE,
    )
    add_textbox(
        slide, emu(1.2), emu(1.95), emu(11.0), emu(0.45),
        "Продавать «Доки» внутри одного разговора вместе с другими сервисами Форус.",
        14, False, WHITE,
    )

    steps = [
        ("1", "Кабинет\nсотрудника", "Вход через HR/кадровый контур"),
        ("2", "Доки", "ЭДО с контрагентами"),
        ("3", "Смартвей", "Командировки и поездки"),
        ("4", "Доки.\nЛогистика", "Перевозочные документы"),
    ]

    left0, top0 = emu(0.97), emu(2.8)
    card_w, card_h, gap = emu(2.7), emu(2.15), emu(0.2)
    for i, (num, title, note) in enumerate(steps):
        left = left0 + i * (card_w + gap)
        fill = BLUE if i in (1, 3) else CARD
        add_card(slide, left, top0, card_w, card_h, fill, 0.08)
        add_textbox(
            slide, left + emu(0.15), top0 + emu(0.2), card_w - emu(0.3), emu(0.4),
            num, 20, True, WHITE, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.1), top0 + emu(0.7), card_w - emu(0.2), emu(0.75),
            title, 14, True, WHITE, PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + emu(0.12), top0 + emu(1.5), card_w - emu(0.24), emu(0.5),
            note, 11, False, WHITE, PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        emu(0.97),
        emu(5.2),
        emu(8.2),
        emu(0.7),
        "Результат: 9 продаж (50% пилота). Скрипт связки работает не хуже точечного канала.",
        14,
        False,
        SOFT,
    )
    add_card(slide, emu(9.4), emu(5.1), emu(3.0), emu(0.9), BLUE, 0.1)
    add_textbox(
        slide, emu(9.5), emu(5.25), emu(2.8), emu(0.6),
        "9 продаж\nпо скрипту", 14, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


def build_slide_6(prs):
    """Сравнение гипотез."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Сравнение двух гипотез", 26)
    clear_body_placeholders(slide)

    rows = [
        ["Параметр", "Гипотеза 1", "Гипотеза 2"],
        ["Суть", "Доки для тех, кому не подходит 1С-ЭПД", "Доки в связке сервисов"],
        ["Вход", "База ЭПД", "Скрипт: КабС → Доки → Смартвей → Доки.Логистика"],
        ["Продажи", "9", "9"],
        ["Доля пилота", "50%", "50%"],
        ["Статус", "Подтверждена", "Подтверждена"],
    ]

    table_shape = slide.shapes.add_table(
        len(rows), 3, emu(0.97), emu(1.65), emu(11.4), emu(3.7)
    )
    table = table_shape.table
    table.columns[0].width = emu(2.4)
    table.columns[1].width = emu(4.5)
    table.columns[2].width = emu(4.5)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)

    add_textbox(
        slide,
        emu(0.97),
        emu(5.55),
        emu(11.4),
        emu(0.6),
        "Вывод: оба канала дают одинаковый результат. Для тиражирования имеет смысл вести оба входа параллельно.",
        13,
        False,
        SOFT,
    )
    return slide


def build_slide_7(prs):
    """Что подготовили."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT2])
    fill_title(slide, "Что уже готово к передаче", 26)
    clear_body_placeholders(slide)

    items = [
        (
            "icon_55.png",
            "Скрипты продаж",
            "Скрипт холодного звонка: Кабинет сотрудника → Доки → Смартвей → Доки.Логистика, "
            "плюс отработка возражений.",
        ),
        (
            "icon_22.png",
            "Материалы для рассылки",
            "Готовые блоки и тексты для партнёрских и клиентских коммуникаций.",
        ),
        (
            "icon_59.png",
            "Сравнение 1С-ЭПД и Доки",
            "Обзор и сравнение: когда предлагать 1С-ЭПД, а когда — Доки / Доки.Логистика.",
        ),
        (
            "icon_67.png",
            "КП, тарифы, шаблоны",
            "Коммерческие предложения, тарифы, клиентские шаблоны презентаций — всё в репозитории.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.6)
    card_w, card_h = emu(5.7), emu(2.2)
    gap_x, gap_y = emu(0.25), emu(0.2)

    for i, (icon_name, title, body) in enumerate(items):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_icon(slide, icon_name, left + emu(0.22), top + emu(0.55), emu(0.7), emu(0.7))
        add_textbox(
            slide, left + emu(1.1), top + emu(0.28), card_w - emu(1.35), emu(0.4),
            title, 15, True, BLUE,
        )
        add_textbox(
            slide, left + emu(1.1), top + emu(0.75), card_w - emu(1.35), emu(1.25),
            body, 12, False, WHITE,
        )
    return slide


def build_slide_8(prs):
    """Оффер."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_BG])
    fill_title(slide, "Оффер, с которым подключали", 26)
    clear_body_placeholders(slide)

    add_card(slide, emu(0.97), emu(1.7), emu(5.5), emu(4.3), BLUE, 0.1)
    add_textbox(
        slide, emu(1.3), emu(2.1), emu(4.8), emu(0.4),
        "Промо на входе", 14, False, WHITE,
    )
    add_textbox(
        slide, emu(1.3), emu(2.6), emu(4.8), emu(1.2),
        "3 месяца\nбезлимит\nбесплатно", 28, True, WHITE,
    )
    add_textbox(
        slide, emu(1.3), emu(4.2), emu(4.8), emu(1.3),
        "Клиент пробует сервис без оплаты пакета. "
        "Барьер входа снимается — проще закрыть первое подключение.",
        14, False, WHITE,
    )

    rights = [
        ("Зачем так делали", "Быстро проверить спрос и набрать практику подключений."),
        ("Что важно дальше", "Отслеживать переход на платный тариф после 3 месяцев."),
        ("Риск", "Без контроля продлений «бесплатный старт» не превратится в выручку."),
        ("Для тиражирования", "Оффер оставить, но добавить KPI по оплате после пробного периода."),
    ]
    y = emu(1.7)
    for title, body in rights:
        add_card(slide, emu(6.7), y, emu(5.65), emu(0.95), CARD, 0.06)
        add_textbox(slide, emu(6.95), y + emu(0.12), emu(5.2), emu(0.28), title, 13, True, BLUE)
        add_textbox(slide, emu(6.95), y + emu(0.45), emu(5.2), emu(0.4), body, 12, False, WHITE)
        y += emu(1.05)
    return slide


def build_slide_9(prs):
    """Выводы."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Выводы по проекту", 26)
    clear_body_placeholders(slide)

    conclusions = [
        (
            "1",
            "Обе гипотезы подтверждены",
            "По 9 продаж на каждый канал. Нет «слабого» направления — оба входа можно масштабировать.",
        ),
        (
            "2",
            "Доки продаётся как отдельно, так и в связке",
            "Работает и как альтернатива 1С-ЭПД, и как часть скрипта с КабС / Смартвей / Доки.Логистика.",
        ),
        (
            "3",
            "Пакет для тиражирования уже собран",
            "Скрипты, рассылки, сравнение, КП и шаблоны лежат в репозитории — новой группе не нужно начинать с нуля.",
        ),
        (
            "4",
            "Пилот закрыл вопрос «продаётся ли»",
            "Следующий этап — не поиск идеи, а системный процесс: объём, конверсия, продление после 3 месяцев.",
        ),
    ]

    y = emu(1.55)
    for num, title, body in conclusions:
        add_card(slide, emu(0.97), y, emu(11.4), emu(1.0), CARD, 0.06)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, emu(1.2), y + emu(0.28), emu(0.45), emu(0.45)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        add_textbox(
            slide, emu(1.2), y + emu(0.28), emu(0.45), emu(0.45),
            num, 14, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(slide, emu(1.9), y + emu(0.15), emu(10.1), emu(0.32), title, 14, True, BLUE)
        add_textbox(slide, emu(1.9), y + emu(0.5), emu(10.1), emu(0.4), body, 13, False, WHITE)
        y += emu(1.15)
    return slide


def build_slide_10(prs):
    """Рекомендации для тиражирования."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Рекомендации для тиражирования", 26)
    clear_body_placeholders(slide)

    recs = [
        (
            "icon_66.png",
            "Два канала сразу",
            "Вести параллельно: (1) база ЭПД → Доки, (2) скрипт связки сервисов.",
        ),
        (
            "icon_54.png",
            "Единый оффер",
            "Сохранить 3 месяца безлимита как стандарт входа, но сразу ставить задачу на продление.",
        ),
        (
            "icon_55.png",
            "Передать пакет as-is",
            "Скрипты, сравнения, рассылки и КП — базовый набор для новой группы без доработки «с нуля».",
        ),
        (
            "icon_64.png",
            "Добавить контроль воронки",
            "Фиксировать: касания → демо → подключения → использование → оплата после пробного периода.",
        ),
    ]

    left0, top0 = emu(0.97), emu(1.6)
    card_w, card_h = emu(5.7), emu(2.2)
    gap_x, gap_y = emu(0.25), emu(0.2)

    for i, (icon_name, title, body) in enumerate(recs):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.08)
        add_icon(slide, icon_name, left + emu(0.22), top + emu(0.55), emu(0.7), emu(0.7))
        add_textbox(
            slide, left + emu(1.1), top + emu(0.28), card_w - emu(1.35), emu(0.4),
            title, 15, True, BLUE,
        )
        add_textbox(
            slide, left + emu(1.1), top + emu(0.75), card_w - emu(1.35), emu(1.25),
            body, 13, False, WHITE,
        )
    return slide


def build_slide_11(prs):
    """Чего не хватает для полной защиты."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Каких данных не хватает для полной защиты", 24)
    clear_body_placeholders(slide)

    add_textbox(
        slide,
        emu(0.97),
        emu(1.45),
        emu(11.4),
        emu(0.45),
        "Ниже — вопросы, которые усилят защиту и передачу. Если данные есть — добавим в презентацию.",
        13,
        False,
        SOFT,
    )

    missing = [
        ("Воронка", "Сколько касаний / звонков / демо на 18 продаж? Какая конверсия по этапам?"),
        ("База", "Размер базы ЭПД и сколько контактов обработали по каждой гипотезе?"),
        ("Деньги", "Средний чек, выручка сейчас и прогноз после окончания 3 бесплатных месяцев?"),
        ("Активность", "Сколько из 18 реально начали отправлять документы?"),
        ("Отказы", "Типовые причины «нет» и где клиенты отваливаются чаще всего?"),
        ("Ресурсы", "Сколько часов менеджеров ушло на пилот? Кто ведёт сопровождение после передачи?"),
    ]

    left0, top0 = emu(0.97), emu(2.0)
    card_w, card_h = emu(5.7), emu(1.25)
    gap_x, gap_y = emu(0.25), emu(0.15)

    for i, (title, body) in enumerate(missing):
        col, row = i % 2, i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, CARD, 0.06)
        add_textbox(slide, left + emu(0.25), top + emu(0.15), card_w - emu(0.5), emu(0.3), title, 13, True, BLUE)
        add_textbox(slide, left + emu(0.25), top + emu(0.5), card_w - emu(0.5), emu(0.6), body, 12, False, WHITE)
    return slide


def build_slide_12(prs):
    """Запрос на решение."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_CONTENT])
    fill_title(slide, "Просим принять решение", 26)
    clear_body_placeholders(slide)

    theses = [
        (
            "Признать пилот успешным",
            "18 подключений за 7 недель, обе гипотезы подтверждены поровну.",
        ),
        (
            "Передать в тиражирование",
            "Новой группе — пакет скриптов, материалов и двух рабочих каналов продаж.",
        ),
        (
            "Закрепить KPI на следующий этап",
            "Подключения + использование + переход на оплату после 3 месяцев.",
        ),
    ]

    y = emu(1.6)
    for i, (title, body) in enumerate(theses):
        add_card(slide, emu(0.97), y, emu(8.2), emu(1.15), CARD, 0.08)
        mark = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, emu(1.2), y + emu(0.35), emu(0.42), emu(0.42)
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = BLUE
        mark.line.fill.background()
        add_textbox(
            slide, emu(1.2), y + emu(0.35), emu(0.42), emu(0.42),
            str(i + 1), 13, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(slide, emu(1.85), y + emu(0.22), emu(7.0), emu(0.35), title, 15, True, BLUE)
        add_textbox(slide, emu(1.85), y + emu(0.6), emu(7.0), emu(0.4), body, 13, False, WHITE)
        y += emu(1.3)

    add_card(slide, emu(9.4), emu(1.6), emu(3.0), emu(4.0), BLUE, 0.1)
    add_icon(slide, "icon_64.png", emu(10.15), emu(1.95), emu(1.5), emu(1.5))
    add_textbox(
        slide, emu(9.6), emu(3.6), emu(2.6), emu(0.9),
        "ГОТОВО\nК ПЕРЕДАЧЕ", 16, True, WHITE, PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, emu(9.6), emu(4.6), emu(2.6), emu(0.7),
        "18 подключений\n2 рабочих канала", 13, False, WHITE, PP_ALIGN.CENTER,
    )
    return slide


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    if not ICONS.exists():
        raise SystemExit(f"Icons not found: {ICONS}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT)

    prs = Presentation(str(OUT))
    delete_all_slides(prs)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
